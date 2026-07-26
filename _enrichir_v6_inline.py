# -*- coding: utf-8 -*-
"""Enrichit le V6 EN INJECTANT directement dans les text frames existants
(blocs « À RETENIR DU BENCHMARK » des slides 4, 5, 7, 8, 9, 10).
Aucun nouveau text box, aucune nouvelle slide. Format inchangé."""
import sys, pathlib, copy
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
from lxml import etree

SRC = pathlib.Path(r'C:\Users\sylva\OneDrive\Bureau\Benchmark_RSE_DRAFTV6_3.pptx')
DST = pathlib.Path(r'C:\Users\sylva\OneDrive\Bureau\Benchmark_RSE_DRAFTV7_enrichi.pptx')

prs = Presentation(str(SRC))

def get_last_para_template(tf):
    """Renvoie le dernier paragraphe pour servir de template (style)."""
    paras = tf.paragraphs
    return paras[-1] if paras else None

def append_para(tf, text, bold=False):
    """Ajoute un paragraphe au text frame en clonant le style du dernier para
    existant (police, taille, couleur). Conserve donc le format."""
    template = get_last_para_template(tf)
    # Crée un nouveau paragraphe en clonant le XML du précédent (tous attributs)
    if template is not None and template._p is not None:
        new_p_xml = deepcopy(template._p)
        # Supprime tous les runs du clone (on les remplace)
        for r in new_p_xml.findall(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}r'):
            new_p_xml.remove(r)
        # Garde les pPr (propriétés de paragraphe) si elles existent
        template._p.addnext(new_p_xml)
        # Wrap en _Paragraph
        from pptx.text.text import _Paragraph
        new_p = _Paragraph(new_p_xml, template._parent)
    else:
        new_p = tf.add_paragraph()
    # Récupère la police/taille/couleur du dernier run de référence
    src_run = None
    if template is not None:
        for r in template.runs:
            src_run = r
    run = new_p.add_run()
    run.text = text
    if src_run is not None:
        try:
            if src_run.font.name:
                run.font.name = src_run.font.name
            if src_run.font.size:
                run.font.size = src_run.font.size
            if src_run.font.color and src_run.font.color.type is not None:
                try:
                    run.font.color.rgb = src_run.font.color.rgb
                except Exception:
                    pass
        except Exception:
            pass
    run.font.bold = bold
    return new_p

def shape_at(slide, idx):
    return list(slide.shapes)[idx]

# ===========================================================
# SLIDE 4 — Ancrage territorial
# ===========================================================
s4 = prs.slides[3]
# Bloc 1 : projets citoyens (shape 15)
tf = shape_at(s4, 15).text_frame
append_para(tf, "GEM, Olympiades de la Sustainability : 15 jours d'activités RSE avec délivrance d'open badges (certificats numériques de soft skills) — valorise l'engagement étudiant sur le CV et LinkedIn.")
# Bloc 2 : compétence étudiants/territoire (shape 20)
tf = shape_at(s4, 20).text_frame
append_para(tf, "Mécénat de compétences : HEC « Stand Up » (programme de 10 semaines pour l'entrepreneuriat social, accessible au personnel) et collaborateurs EMLYON autorisés à consacrer du temps de travail à des associations partenaires.")
# Bloc 3 : achat local (shape 25) — déjà très riche, on n'ajoute rien

# ===========================================================
# SLIDE 5 — Engager parties prenantes
# ===========================================================
s5 = prs.slides[4]
# Bloc 1 : former enseignants (shape 15) — déjà riche
# Bloc 2 : référents/ambassadeurs (shape 20) — déjà riche
# Bloc 3 : apprenants vs collab (shape 25)
tf = shape_at(s5, 25).text_frame
append_para(tf, "LABELS D'ENGAGEMENT  EDHEC label « Citizen Associations » (cahier des charges 3 critères, objectif 100 % des assos labellisées d'ici 2028) ; EMLYON charte RSE signée par 42 associations + label événements Ecollab (64 % obtenu).")

# ===========================================================
# SLIDE 7 — Environnement (alimentation, numérique, mobilité)
# ===========================================================
s7 = prs.slides[6]
# Bloc 2 : numérique responsable (shape 20)
tf = shape_at(s7, 20).text_frame
append_para(tf, "Quantifier le CO2 évité : HEC publie 5,9 tCO2 évitées par le reconditionnement IT et 2,2 tCO2 par le remplacement de batteries ; recettes de revente reversées à un Sustainability Fund interne.")

# ===========================================================
# SLIDE 8 — Environnement (biodiversité, déchets, achats)
# ===========================================================
s8 = prs.slides[7]
# Bloc 3 : achats responsables (shape 25)
tf = shape_at(s8, 25).text_frame
append_para(tf, "Verrouiller via les RH : AD Education indexe une part des primes de son DG sur les critères ESG — pratique très peu commune mais signal de gouvernance le plus fort du benchmark (1 entité).")

# ===========================================================
# SLIDE 9 — Gouvernance responsable
# ===========================================================
s9 = prs.slides[8]
# Bloc 1 : gouvernance RSE (shape 15)
tf = shape_at(s9, 15).text_frame
append_para(tf, "Faire vérifier par un tiers indépendant : GEM fait auditer ses indicateurs RSE par KPMG ; EMLYON est accréditée COFRAC pour son Comité de Mission et a obtenu le label européen HRS4R (Human Resource Strategy for Researchers, Commission Européenne, fév. 2025).")
# Bloc 2 : compliance/éthique (shape 20)
tf = shape_at(s9, 20).text_frame
append_para(tf, "Rémunération indexée RSE : AD Education indexe une part des primes de son DG sur les critères ESG (très peu pratiqué, sujet de gouvernance émergent et signal le plus fort observé).")

# ===========================================================
# SLIDE 10 — QVT & égalité des chances
# ===========================================================
s10 = prs.slides[9]
# Bloc 1 : VSS (shape 15) — on ajoute la santé mentale (Nightline, sentinelles, VR)
tf = shape_at(s10, 15).text_frame
append_para(tf, "SANTÉ MENTALE  Écoute par les pairs : EDHEC et EMLYON ont déployé Nightline (écoute nocturne gratuite par des étudiants formés). Sentinelles formées chez EDHEC, et chez ESSEC enseignants, résidences et bureaux d'associations sont formés au repérage. Hotline 24/7 externalisée (Pros-Consulte chez EDHEC, APICIL chez EMLYON). KEDGE a ouvert en oct. 2025 la 1ère Maison de la Santé et du Bien-Être pour une école en France.")
append_para(tf, "INNOVATION  EMLYON organise des Journées de l'inclusion en réalité virtuelle (startup Reverto qu'elle a incubée) : on vit une situation de handicap ou une scène de sexisme — peu commun, mémorable, transposable dans HOPEN.")
# Bloc 2 : inclusion / égalité chances (shape 20)
tf = shape_at(s10, 20).text_frame
append_para(tf, "Mixité publiée en détail (exemplarité interne) : EMLYON affiche COMEX 40 % de femmes, CoDir 47 % (+18 % depuis 2021), Conseil de Surveillance 41 %, corps professoral 37 % ; EDHEC suit le % de présidentes d'associations (47 %). Index F/H détaillé poste par poste chez EMLYON (rémunération, augmentations, promotions, retour congé maternité).")
# Bloc 4 : fonds de dotation, financer quoi (shape 25 ou plus loin)
# La slide 10 a 4 blocs au lieu de 3, le bloc 4 est shape suivante
# Recherchons le bloc 4 dynamiquement (texte contient "fonds de dotation")
fond_block = None
for sh in s10.shapes:
    if sh.has_text_frame and 'fonds' in sh.text_frame.text.lower() and 'urgence' in sh.text_frame.text.lower():
        fond_block = sh
        break
if fond_block is not None:
    tf = fond_block.text_frame
    append_para(tf, "BOURSES CIBLÉES DISTINCTIVES  Bourses réfugiés : HEC « Imagine Fellows » (zones de conflit), BSB « (Re)Connect », GEM Refugee Grant Program, PSB « Each One ». Prêt étudiant co-garanti école + banque : KEDGE KAP Equal Opportunities (50 % école / 50 % banque) — lève la barrière du financement initial sans dépense directe. Bourse mobilité bas carbone (train vs avion) : ESSEC, EFREI (61 bourses), EMLYON, TSM.")
else:
    print('AVERTISSEMENT : bloc « fonds de dotation » non trouvé en slide 10')

prs.save(str(DST))
print(f'✓ Sauvegardé : {DST}')
print(f'Slides : {len(prs.slides)} (inchangé — pas de nouvelle slide ajoutée)')

# -*- coding: utf-8 -*-
"""Inspecte la structure du V6 pour identifier où injecter le contenu."""
import sys, pathlib
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Emu

SRC = pathlib.Path(r'C:\Users\sylva\OneDrive\Bureau\Benchmark_RSE_DRAFTV6_3.pptx')
prs = Presentation(str(SRC))

def emu_to_in(e):
    return e / 914400 if e else 0

# Inspecte les slides 4 à 10 (index 3 à 9)
TARGETS = [3, 4, 6, 7, 8, 9]  # 0-based : 4=Ancrage, 5=Engager, 7=Env1, 8=Env2, 9=Gouv, 10=QVT
for idx in TARGETS:
    s = prs.slides[idx]
    print(f'\n========= SLIDE {idx + 1} =========')
    for i, sh in enumerate(s.shapes):
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text.strip()
            x, y = emu_to_in(sh.left), emu_to_in(sh.top)
            w, h = emu_to_in(sh.width), emu_to_in(sh.height)
            preview = txt[:120].replace('\n', ' | ')
            print(f'  [{i:2}] x={x:5.2f} y={y:5.2f} w={w:5.2f} h={h:5.2f} '
                  f'paras={len(sh.text_frame.paragraphs):2}  "{preview}"')

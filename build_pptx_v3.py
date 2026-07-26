# -*- coding: utf-8 -*-
"""V3 — 2 slides par rapport (visuels + analyse approfondie Forme/Fond/Idées concrètes)."""
import sys, json, pathlib, re
sys.stdout.reconfigure(encoding='utf-8')

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ===== Charte =====
ROSE = RGBColor(0xE6, 0x0F, 0x7D)
ROSE_LIGHT = RGBColor(0xFD, 0xE3, 0xF1)
VIOLET = RGBColor(0x4A, 0x19, 0x42)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xFA, 0xFA, 0xFA)
TEXT = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
RED = RGBColor(0xE5, 0x39, 0x35)
RED_BG = RGBColor(0xFF, 0xEB, 0xEE)
YELLOW = RGBColor(0xFF, 0xC1, 0x07)
YELLOW_BG = RGBColor(0xFF, 0xF8, 0xE1)
GREY = RGBColor(0xEC, 0xEC, 0xEC)

LOGO_PATH = pathlib.Path('logo-igensia.png')
COVERS_DIR = pathlib.Path('_covers')
META_PATH = COVERS_DIR / '_meta.json'

# ===== Données ranking =====
raw = pathlib.Path('data.js').read_text(encoding='utf-8')
m = re.match(r'^(\s*const\s+BENCHMARK_DATA\s*=\s*)(.+?)(;\s*)$', raw, re.DOTALL)
data = json.loads(m.group(2))

def parse(v):
    try:
        n = float(v); return n if 0 <= n <= 5 else None
    except: return None
def calc_final(notes):
    vals = []
    for who in ['blanche','sylvain']:
        for typ in ['forme','fond']:
            n = parse(notes.get(who,{}).get(typ,''))
            if n is not None: vals.append(n)
    return sum(vals)/len(vals) if vals else None

focus_data = {}
for f in data['focus']:
    nm = f['name']
    notes = f['notes']
    bf = parse(notes['blanche'].get('forme',''))
    bb = parse(notes['blanche'].get('fond',''))
    sf = parse(notes['sylvain'].get('forme',''))
    sb = parse(notes['sylvain'].get('fond',''))
    avg_f = (bf+sf)/2 if (bf is not None and sf is not None) else None
    avg_b = (bb+sb)/2 if (bb is not None and sb is not None) else None
    fn = calc_final(notes)
    short = nm.split('\n')[0].replace('★ ','').replace('BUSINESS SCHOOL','BS').strip()
    focus_data[nm] = {
        'short': short, 'avg_f': avg_f, 'avg_b': avg_b, 'final': fn,
        'is_igensia': 'IGENSIA' in nm.upper(),
    }
ranking = [(nm, d) for nm, d in focus_data.items()]
ranking.sort(key=lambda x: -x[1]['final'] if x[1]['final'] is not None else 0)
for i, (nm, d) in enumerate(ranking):
    d['rank'] = i + 1
igensia_rank = next(d['rank'] for nm, d in ranking if d['is_igensia'])

COVER_KEYS = {
    '★ IGENSIA EDUCATION': 'igensia', 'OMNES EDUCATION': 'omnes',
    'GALILEO GLOBAL EDUCATION': 'galileo', 'AD EDUCATION': 'ad',
    'ESSEC BUSINESS SCHOOL': 'essec', 'EMLYON BUSINESS SCHOOL': 'emlyon',
    'EDHEC BUSINESS SCHOOL': 'edhec', 'INSEAD': 'insead', 'HEC PARIS': 'hec',
    'EXCELIA BUSINESS SCHOOL': 'excelia', 'AUDENCIA BUSINESS SCHOOL': 'audencia',
    'KEDGE BUSINESS SCHOOL': 'kedge', 'SKEMA BUSINESS SCHOOL': 'skema',
    'GRENOBLE EM': 'gem', 'TBS': 'tbs', 'IESEG SCHOOL OF MANAGEMENT': 'ieseg',
    'NEOMA BUSINESS SCHOOL': 'neoma', 'BURGUNDY SCHOOL OF BUSINESS': 'bsb',
    'EFREI': 'efrei', 'EM STRASBOURG BUSINESS SCHOOL': 'emstras',
    'PARIS SCHOOL OF BUSINESS\n(Groupe Galileo)': 'psb', 'CESI': 'cesi',
}

meta = {}
if META_PATH.exists():
    try: meta = json.loads(META_PATH.read_text(encoding='utf-8'))
    except: pass

# ===== ANALYSES — adaptées au contexte école/groupe ES =====
# Pour chaque école : titre rapport, pages, légendes des images, forme +/-, fond +/-, idées concrètes
ANALYSES = {
    '★ IGENSIA EDUCATION': {
        'titre': 'Rapport RSE 2024-2025',
        'pages': '48 p.',
        'leg_cov': "Couverture violet IGENSIA + rose fuchsia, photo apprenants en gros plan",
        'leg_i1':  "Sommaire en 5 axes, hiérarchie visuelle claire (photos pleine page)",
        'leg_i2':  "Page chiffres-clés avec codification « C'est réalisé / lancé / prévu »",
        'forme_plus': [
            "Identité visuelle forte et cohérente (violet + rose, photos qualitative en pleine page)",
            "Codification claire « C'est réalisé / lancé / prévu » qui signale l'avancement de chaque action",
        ],
        'forme_moins': [
            "Densité photographique élevée qui prend de la place au détriment des KPIs",
            "Pas de témoignages incarnés (collaborateurs, apprenants, alumni)",
        ],
        'fond_plus': [
            "Distinction apprenants/collaborateurs intégrée dans chaque section (vraie singularité)",
            "Bilan carbone détaillé : 25 498 tCO2e ventilé par poste et par campus",
            "Couverture complète des 5 axes du référentiel LUCIE 26000",
        ],
        'fond_moins': [
            "Index égalité F/H en recul (84 → 75) non commenté",
            "Dispositif VSS pas encore déployé",
            "Budget RSE non publié (pas d'enveloppe chiffrée)",
        ],
        'idees': [
            ("C'est notre point de référence", "Ce rapport est la base à améliorer pour le prochain cycle"),
        ],
    },
    'OMNES EDUCATION': {
        'titre': 'Rapport RSE 2024-2025',
        'pages': '44 p.',
        'leg_cov': "Couverture sobre, photo apprenants + accroche Time to Act",
        'leg_i1':  "Page sommaire : 4 grands axes avec icônes + photos d'illustration",
        'leg_i2':  "Tableau d'évolution du budget RSE (310k€) depuis 2020",
        'forme_plus': [
            "Mise en page A4 paysage en double-page : lecture continue, fil narratif fort",
            "Iconographie cohérente et reconnaissable (« Time to Act »)",
            "Rythme visuel équilibré : photo / texte / KPI / témoignage",
        ],
        'forme_moins': [
            "Format paysage moins pratique à lire sur écran qu'A4 portrait",
            "Quelques pages très denses en texte (axes 3 et 4)",
        ],
        'fond_plus': [
            "Budget RSE publié 310k€ avec son évolution depuis 2020 (acte de transparence rare)",
            "Bilan carbone scopes 1+2+3 sur 17 720 tCO2e (506 kgCO2e/étudiant)",
            "Mention explicite des objectifs non atteints (rare et crédibilisant)",
            "Distinction apprenants/collaborateurs explicite tout au long du document",
        ],
        'fond_moins': [
            "Peu de témoignages incarnés (apprenants, collaborateurs, partenaires)",
            "Plan d'action 2025-2030 présenté de façon textuelle, mériterait un schéma",
        ],
        'idees': [
            ("Publier notre budget RSE avec son évolution sur 4-5 ans",
             "Pourquoi : c'est le signal de transparence le plus fort. Comment : 1 tableau simple budget total + ETP + % budget global, comparé année par année."),
            ("Lister les objectifs non atteints",
             "Pourquoi : ça distingue un rapport d'engagement d'un document de communication. Comment : à la fin de chaque axe, une mention « ce que nous n'avons pas atteint en 2024 » + explication courte."),
        ],
    },
    'GALILEO GLOBAL EDUCATION': {
        'titre': "Rapport d'Impact 2024-2025",
        'pages': '122 p.',
        'leg_cov': "Couverture moderne, photo aérienne + accroche « Forward Impact »",
        'leg_i1':  "Page matrice de matérialité (double matérialité visuelle)",
        'leg_i2':  "Mosaïque d'écoles : présentation visuelle des 56 marques du Groupe",
        'forme_plus': [
            "Identité graphique très soignée, niveau corporate international",
            "Mosaïques visuelles efficaces pour présenter les 56 écoles du Groupe",
            "Matrice de matérialité bien posée graphiquement",
        ],
        'forme_moins': [
            "122 pages : volume écrasant, le lecteur décroche",
            "KPIs noyés dans la masse, pas de page de synthèse en début",
            "Frontière floue entre « actions réalisées » et « actions prévues »",
        ],
        'fond_plus': [
            "Couverture internationale (110 000 étudiants, 18 pays) avec exemples par marque",
            "Matrice de matérialité comme outil d'analyse stratégique",
        ],
        'fond_moins': [
            "Bilan carbone limité aux campus, scope 3 absent",
            "Manque de hiérarchie : tout semble équivalent",
        ],
        'idees': [
            ("Construire une matrice 2×2 « priorité × impact » des enjeux RSE",
             "Pourquoi : permet au lecteur de comprendre pourquoi on agit sur tel sujet plutôt qu'un autre. Comment : 1 page en début de rapport, 8-10 enjeux placés sur 2 axes."),
            ("Ne pas dépasser 50 pages",
             "Pourquoi : au-delà, le rapport devient illisible. Comment : 1 page de synthèse par axe + renvois vers annexes web pour le détail."),
        ],
    },
    'AD EDUCATION': {
        'titre': 'Rapport ESG 2024',
        'pages': '34 p.',
        'leg_cov': "Couverture corporate, codes financiers (gris/bleu)",
        'leg_i1':  "Tableau d'indicateurs ESG : KPIs alignés sur normes investisseurs",
        'leg_i2':  "Page « engagements » avec icônes ODD",
        'forme_plus': [
            "Format A4 sobre, lisibilité sans fioritures",
            "Tableaux d'indicateurs propres, faciles à lire",
        ],
        'forme_moins': [
            "Codes très financiers (gris/bleu) qui n'engagent pas le lecteur étudiant/collaborateur",
            "Pas de photos, pas de témoignages, pas de mise en récit",
        ],
        'fond_plus': [
            "Indicateurs ESG normés (utile pour partenaires investisseurs)",
        ],
        'fond_moins': [
            "Bilan carbone limité aux activités hors France (2 000 tCO2e seulement) : périmètre douteux",
            "Certification « Neutre Carbone » limitée aux scopes 1+2 (méthodologie discutable)",
            "Axes gouvernance et recherche à peine traités",
            "Pas de distinction apprenants/collaborateurs",
        ],
        'idees': [
            ("Contre-exemple : éviter le format ESG investisseurs",
             "Pourquoi : un rapport RSE école s'adresse d'abord à ses apprenants, collaborateurs, alumni, candidats. Pas à des analystes financiers. Comment : privilégier un ton incarné, des photos et des témoignages."),
        ],
    },
    'ESSEC BUSINESS SCHOOL': {
        'titre': 'Rapport DD&RS 2025',
        'pages': '22 p.',
        'leg_cov': "Couverture sobre, photo apprenants + titre épuré",
        'leg_i1':  "Schéma directeur DD&RS 2024 : 1 page avec les 5 axes et leurs priorités",
        'leg_i2':  "Page « Réalisé / Nouveau / En cours » sur 1 axe particulier",
        'forme_plus': [
            "Densité éditoriale rare : 22 pages très denses sans être indigestes",
            "Structure « C'est réalisé / C'est nouveau / C'est en cours » par axe (très lisible)",
            "Schéma directeur visible et lisible dès l'entrée du rapport",
        ],
        'forme_moins': [
            "Photos peu nombreuses, ton un peu froid",
            "Peu de témoignages",
        ],
        'fond_plus': [
            "37 ETP dédiés à la RSE annoncés clairement",
            "130 heures de cours sur la transition (volume horaire élevé documenté)",
            "Bilan carbone complet + matrice de matérialité",
        ],
        'fond_moins': [
            "Ton trop maîtrisé : peu d'autocritique, peu de signal sur les difficultés",
            "Mentionne peu les apprenants comme parties prenantes actives",
        ],
        'idees': [
            ("Faire un schéma directeur visible dès l'entrée du rapport",
             "Pourquoi : permet au lecteur de comprendre où on va avant d'entrer dans le détail. Comment : 1 page avec les 5 axes IGENSIA + 2-3 priorités par axe, sur une trajectoire 2024-2028."),
            ("Adopter la structure « Réalisé / Nouveau / En cours » par axe",
             "Pourquoi : c'est proche de notre codification actuelle mais plus structurée. Comment : pour chaque axe, 3 sous-titres fixes avec puces sous chacun. Cohérence éditoriale sur l'ensemble du rapport."),
        ],
    },
    'EMLYON BUSINESS SCHOOL': {
        'titre': "Rapport d'engagement 2024",
        'pages': '38 p.',
        'leg_cov': "Couverture engageante, accroche « Société à Mission »",
        'leg_i1':  "Page « SDGs Inside » : tableau de couverture des cours par ODD",
        'leg_i2':  "Témoignages d'apprenants en double-page",
        'forme_plus': [
            "Mise en page dynamique, beaucoup de photos pleine page",
            "Page « SDGs Inside » : un schéma fort qui synthétise la couverture des cours",
            "Témoignages d'apprenants en double-page : incarnation forte",
        ],
        'forme_moins': [
            "Ton très promotionnel, peu de recul critique",
            "Mise en page parfois chargée (3-4 informations différentes par page)",
        ],
        'fond_plus': [
            "Référentiel SDGs Inside : 100% des cours analysés sous l'angle ODD",
            "Bilan carbone affichant -30%",
            "Statut Société à Mission utilisé comme fil rouge",
        ],
        'fond_moins': [
            "Politique sociale collaborateurs peu développée",
            "Limites et difficultés peu présentes",
        ],
        'idees': [
            ("Analyser nos cours sous l'angle ODD (« SDG Mapping »)",
             "Pourquoi : objective notre contribution pédagogique à la RSE. Comment : tableau « ODD × cours du tronc commun » sur 1 page, lisible d'un coup d'œil."),
            ("Intégrer 4-6 témoignages courts d'apprenants",
             "Pourquoi : ça humanise le rapport et donne des « preuves de vie ». Comment : photos qualitatives + verbatim de 3-4 lignes + mention du programme/promo."),
        ],
    },
    'EDHEC BUSINESS SCHOOL': {
        'titre': 'Rapport DDRS 2023',
        'pages': '36 p.',
        'leg_cov': "Couverture artistique, illustration sur fond bleu marine",
        'leg_i1':  "Campagne anti-VSS : visuel fort + chiffres-clés",
        'leg_i2':  "Matrice de matérialité (Capgemini Invent)",
        'forme_plus': [
            "Illustrations très travaillées tout au long du rapport (qualité graphique haute)",
            "Campagne anti-VSS particulièrement réussie visuellement",
            "Rythme : 1 enjeu = 1 double-page (illustration gauche + texte droite)",
        ],
        'forme_moins': [
            "Couleurs très sombres (bleu marine dominant) qui peuvent fatiguer la lecture",
            "Texte parfois trop dense par bloc",
        ],
        'fond_plus': [
            "Matrice de matérialité réalisée avec Capgemini Invent (audit méthodologique)",
            "Bilan carbone complet, 100% étudiants sensibilisés",
        ],
        'fond_moins': [
            "Environnement campus moins détaillé que les autres axes",
            "Pas d'audit externe du bilan carbone",
        ],
        'idees': [
            ("Soigner le traitement graphique des sujets sensibles (VSS, handicap, diversité)",
             "Pourquoi : un visuel fort traduit la considération qu'on porte au sujet. Comment : commander à un graphiste 2-3 illustrations originales pour les axes sensibles."),
            ("Rythme « 1 enjeu = 1 double-page »",
             "Pourquoi : crée un fil narratif clair, le lecteur sait toujours où il en est. Comment : structurer 5 axes × N enjeux = N×2 pages, chaque double-page commence par une illustration ou photo."),
        ],
    },
    'INSEAD': {
        'titre': 'Sustainability Report 2023',
        'pages': '67 p.',
        'leg_cov': "Couverture institutionnelle, photo campus + titre",
        'leg_i1':  "Page académique : présentation des 35 cas pédagogiques RSE",
        'leg_i2':  "Bilan carbone détaillé par campus (graphiques)",
        'forme_plus': [
            "Bilan carbone détaillé par campus : transparence sur la dispersion géographique",
            "Section recherche très structurée",
        ],
        'forme_moins': [
            "Très papier de recherche : texte dense, peu de visuels, pas de témoignages",
            "67 pages : volume excessif pour un public général",
            "Hiérarchie typographique faible (peu de respirations)",
        ],
        'fond_plus': [
            "35 cas pédagogiques + 56 articles RSE : production académique remarquable",
            "Bilan carbone détaillé par campus (Fontainebleau, Singapour, San Francisco)",
        ],
        'fond_moins': [
            "Actions concrètes sur les campus peu présentes",
            "Volet social et engagement étudiant en retrait",
            "Ton institutionnel, peu incarné",
        ],
        'idees': [
            ("Valoriser la production académique de nos formateurs",
             "Pourquoi : prouve que notre engagement RSE n'est pas que de la communication. Comment : page « La recherche RSE chez nous » avec liste des publications, des interventions et des projets de recherche en lien."),
        ],
    },
    'HEC PARIS': {
        'titre': 'Sustainability Report 2024',
        'pages': '52 p.',
        'leg_cov': "Couverture sobre, photo campus + titre épuré",
        'leg_i1':  "Page « Audited Carbon Footprint » : signature AFNOR Certification",
        'leg_i2':  "Page recherche : 20% des publications portent sur la durabilité",
        'forme_plus': [
            "Format propre, lecture confortable, hiérarchie typographique soignée",
            "Mise en avant visuelle de l'audit AFNOR du bilan carbone",
        ],
        'forme_moins': [
            "Très institutionnel, peu de mise en récit",
            "Peu de témoignages, peu de photos d'actions de terrain",
            "Données opérationnelles sur le quotidien des campus peu présentes",
        ],
        'fond_plus': [
            "Bilan carbone audité par AFNOR Certification (rare et fort)",
            "20% de la recherche porte sur la durabilité (chiffre fort)",
        ],
        'fond_moins': [
            "Couverture des 5 axes correcte sans aller au bout sur chacun",
            "Engagement étudiant peu détaillé",
        ],
        'idees': [
            ("Faire auditer notre bilan carbone par un tiers",
             "Pourquoi : c'est le signal de fiabilité le plus puissant. Pour 25 498 tCO2e, ça crédibilise tout le rapport. Comment : prestataire (AFNOR / EY / Bureau Veritas), coût ~10-15k€, à prévoir en année N-1."),
        ],
    },
    'EXCELIA BUSINESS SCHOOL': {
        'titre': 'Rapport TES 2024',
        'pages': '44 p.',
        'leg_cov': "Couverture engageante, mise en avant de l'IRSI (Institut RSE)",
        'leg_i1':  "Page chiffres-clés : 28 projets, 13 903h DD, etc.",
        'leg_i2':  "Tableau d'avancement par axe avec mention des objectifs non atteints",
        'forme_plus': [
            "Distinction apprenants/collaborateurs présente dans la structure des sections",
            "Tableaux d'avancement par axe : visualisation claire de l'état d'avancement",
        ],
        'forme_moins': [
            "Format dense, peu de respirations visuelles entre les chiffres",
            "Quelques pages trop chargées en chiffres",
        ],
        'fond_plus': [
            "Historique long (IRSI depuis 2010) qui crédibilise la démarche",
            "KPIs nombreux et précis : 28 projets transversaux, 13 903h DD",
            "Mention des objectifs non atteints (signal de maturité)",
        ],
        'fond_moins': [
            "Format dense rend la lecture rapide difficile",
        ],
        'idees': [
            ("Multiplier les KPIs chiffrés (volume horaire, nb projets, % couverture)",
             "Pourquoi : un rapport sans KPIs reste déclaratif. Comment : pour chaque axe, 3-4 indicateurs simples (nb d'apprenants concernés, h de formation, % de couverture des programmes)."),
            ("Mettre en avant l'historique d'engagement (« X années »)",
             "Pourquoi : crédibilise la démarche. Comment : timeline en début de rapport « Notre engagement RSE depuis YYYY » avec 5-6 jalons."),
        ],
    },
    'AUDENCIA BUSINESS SCHOOL': {
        'titre': 'Rapport DD&RS 2026',
        'pages': '16 p.',
        'leg_cov': "Couverture compacte, accroche « Pionnière Label DD&RS depuis 2016 »",
        'leg_i1':  "Sommaire compact 4 axes + 5 priorités stratégiques",
        'leg_i2':  "Tableau comparatif 2024/2025 : 15 KPIs avec évolution chiffrée (page 15)",
        'forme_plus': [
            "Format compact (16 pages) sans sacrifier la profondeur : densité éditoriale exemplaire",
            "Structure « Réalisations + Projections + Chiffres clés » très lisible",
            "Le tableau comparatif p.15 (15 KPIs en 1 page) est la vraie réussite éditoriale",
        ],
        'forme_moins': [
            "Format court limite la profondeur d'analyse par axe",
            "Pas de témoignages",
        ],
        'fond_plus': [
            "Pionnière 1ère BS Label DD&RS (2016) : argument d'autorité",
            "PRME Champions (2014), Global Compact (2004) : 20 ans d'historique",
            "Tableau d'évolution chiffrée d'une année sur l'autre",
            "Signaux honnêtes : budget transition en baisse (11→8%), CO2/étudiant en hausse",
        ],
        'fond_moins': [
            "Format court oblige à des choix",
        ],
        'idees': [
            ("Tableau comparatif annuel des 12-15 KPIs",
             "Pourquoi : c'est l'élément qui transforme un état des lieux en trajectoire. Comment : à la fin du rapport, 1 page avec 12-15 indicateurs en lignes, colonnes N-1 et N, et flèche évolution. À reprendre dans le rapport de l'année suivante."),
            ("Structure « Réalisations / Projections / Chiffres-clés » par axe",
             "Pourquoi : structure répétée = lecture rapide. Comment : pour chaque axe, 3 sous-blocs fixes."),
        ],
    },
    'KEDGE BUSINESS SCHOOL': {
        'titre': 'Rapport DD 2020-2021',
        'pages': '16 p.',
        'leg_cov': "Couverture compacte avec rappel des engagements",
        'leg_i1':  "Page engagements : 5 axes (Pioneering / Champion / etc.)",
        'leg_i2':  "Page KPIs : 150k€ budget, 70 volontaires",
        'forme_plus': [
            "Format compact, lecture rapide",
            "Identité visuelle KEDGE bleue/violet cohérente",
        ],
        'forme_moins': [
            "Rapport ancien (2020-2021), aspect visuellement daté",
            "Densité d'information faible par page",
        ],
        'fond_plus': [
            "KPIs présents (150k€ budget, 70 volontaires)",
        ],
        'fond_moins': [
            "Rapport ancien : ne reflète plus les actions récentes",
            "Bilan carbone scope 3 mentionné sans détail",
            "5 ans sans publication : signal négatif",
        ],
        'idees': [
            ("Ne pas attendre 4-5 ans entre 2 rapports",
             "Pourquoi : un rapport ancien finit par décrédibiliser l'engagement. Comment : rythme annuel minimum, version « light » possible (16 pages) en année creuse."),
        ],
    },
    'SKEMA BUSINESS SCHOOL': {
        'titre': 'Rapport Transitions 2024-2025',
        'pages': '24 p.',
        'leg_cov': "Couverture moderne, accroche « Transitions Act »",
        'leg_i1':  "Page Fresque du Climat : 1 300 participants chiffré",
        'leg_i2':  "Page labels (DD&RS, ISO 14001, Charte INR)",
        'forme_plus': [
            "Identité graphique moderne, lecture confortable",
            "Mise en avant des labels (DD&RS, ISO 14001, INR) en début de rapport",
        ],
        'forme_moins': [
            "Mélange entre actions « réalisé » et « prévu » pas toujours clair",
            "24 pages mais peu de détails par axe",
        ],
        'fond_plus': [
            "Trio de labels (DD&RS + ISO 14001 + Charte INR) : signal externe puissant",
            "Fresque du Climat à 1 300 participants : preuve de déploiement",
        ],
        'fond_moins': [
            "Détail léger pour une école labellisée DD&RS (déçoit l'attente)",
            "Mélange réalisé/prévu pas toujours lisible",
        ],
        'idees': [
            ("Combiner plusieurs labels externes en signature de bas de page",
             "Pourquoi : un seul label peut paraître anecdotique, plusieurs ensemble crédibilisent. Comment : bandeau « Nos engagements externes : LUCIE 26000 + ... » à présent dès la page 2."),
        ],
    },
    'GRENOBLE EM': {
        'titre': "Communication sur l'Engagement 2025",
        'pages': '38 p.',
        'leg_cov': "Couverture statut Société à Mission, photo campus",
        'leg_i1':  "Trajectoire SBTi validée : graphique de réduction GES",
        'leg_i2':  "Comité de Mission indépendant : composition + rôle",
        'forme_plus': [
            "Graphique de trajectoire SBTi très efficace visuellement",
            "Format sobre mais bien structuré",
        ],
        'forme_moins': [
            "Plus institutionnel que graphique",
            "Peu de témoignages incarnés",
        ],
        'fond_plus': [
            "Bilan carbone 10 314 tCO2e (-14% vs 2019) avec trajectoire SBTi validée",
            "Comité de Mission indépendant : signal de gouvernance fort",
            "Statut Société à Mission (1ère grande école française, 2021)",
        ],
        'fond_moins': [
            "Engagement éditorial qui n'engage pas le lecteur",
        ],
        'idees': [
            ("Visualiser la trajectoire de réduction GES sous forme de graphique",
             "Pourquoi : un graphique de trajectoire vaut mille mots. Comment : courbe avec point de départ (2024 : 25 498 tCO2e), jalons annuels et objectif 2030. Même si on n'a pas SBTi, une trajectoire visuelle a déjà un impact."),
            ("Faire valider notre engagement par un comité externe ou un tiers",
             "Pourquoi : crédibilise la démarche. Comment : peut-être pas un comité de mission complet, mais une lettre d'évaluation annuelle par un tiers reconnu."),
        ],
    },
    'TBS': {
        'titre': 'Rapport de Société à Mission 2024-2025',
        'pages': '30 p.',
        'leg_cov': "Couverture format paysage 16:9 (peu courant)",
        'leg_i1':  "Page recherche : 42 publications RSE dans des revues 3-4 étoiles",
        'leg_i2':  "Présentation des 4 objectifs statutaires avec indicateurs",
        'forme_plus': [
            "Format paysage 16:9 original (lisible à l'écran)",
            "Page de recherche très claire (publications quantifiées)",
        ],
        'forme_moins': [
            "Format paysage moins pratique à imprimer",
            "Ton institutionnel, peu de mise en récit",
            "Peu de photos qualitatives",
        ],
        'fond_plus': [
            "61% des publications portent sur la RSE, dont 42 dans des revues 3-4 étoiles",
            "4 objectifs statutaires (Société à Mission) suivis avec indicateurs",
            "Comité de Mission documenté",
        ],
        'fond_moins': [
            "Environnement campus peu couvert (axes équipements, déchets, biodiversité absents)",
        ],
        'idees': [
            ("Valoriser les travaux RSE de nos formateurs (publications, conférences, projets de recherche)",
             "Pourquoi : prouve que notre engagement va au-delà du déclaratif. Comment : page « Notre matière grise au service de la RSE » avec liste des publications, conférences, projets en lien."),
        ],
    },
    'IESEG SCHOOL OF MANAGEMENT': {
        'titre': "Rapport d'Impact 2024-2025",
        'pages': '47 p.',
        'leg_cov': "Couverture format paysage moderne",
        'leg_i1':  "Tableau « 81% des cours du PGE intègrent la durabilité »",
        'leg_i2':  "Page « 61 cours entièrement dédiés à la durabilité » avec liste",
        'forme_plus': [
            "Structure en 4 piliers claire",
            "Bilan carbone publié directement dans le rapport (pas en annexe)",
        ],
        'forme_moins': [
            "Structure en 4 piliers un peu classique, peu d'originalité éditoriale",
            "Format paysage qui rend la mise en page moins flexible",
        ],
        'fond_plus': [
            "Intégration pédagogique très forte : 81% des cours du PGE intègrent la durabilité",
            "61 cours entièrement dédiés à la durabilité (liste accessible)",
            "Distinction apprenants/collaborateurs claire",
        ],
        'fond_moins': [
            "Manque d'autocritique sur les axes moins avancés",
        ],
        'idees': [
            ("Quantifier le % de cours intégrant un module DD/RSE",
             "Pourquoi : la pédagogie est notre cœur de métier, c'est le KPI le plus naturel pour nous. Comment : audit des syllabi de nos 11 écoles, classification simple « Module DD/RSE oui/non ». Publier le %."),
            ("Publier la liste des cours/modules spécifiquement dédiés à la RSE/DD",
             "Pourquoi : ça matérialise notre engagement pédagogique. Comment : 1 page (ou annexe web liée) avec la liste des cours, par école, par niveau."),
        ],
    },
    'NEOMA BUSINESS SCHOOL': {
        'titre': "Rapport d'Engagement TSE 2025",
        'pages': '44 p.',
        'leg_cov': "Couverture engageante avec photo apprenants",
        'leg_i1':  "Index F/H : 88/100 + 52% étudiantes (chiffres mis en avant)",
        'leg_i2':  "Page dispositif VSS : adresse angela@neoma-bs.fr en gros",
        'forme_plus': [
            "Adresse VSS « angela@neoma-bs.fr » mise en avant graphiquement (1 page entière)",
            "Photos apprenants en double-page",
        ],
        'forme_moins': [
            "Ton très promotionnel",
            "Mise en perspective des difficultés absente",
        ],
        'fond_plus': [
            "Index égalité F/H 88, 52% étudiantes (vs IGENSIA 75)",
            "Adresse dédiée pour les signalements VSS : dispositif visible et accessible",
            "Bilan carbone scopes 1-3",
        ],
        'fond_moins': [
            "Recherche peu détaillée",
            "Limites et difficultés peu présentes",
        ],
        'idees': [
            ("Mettre en place une adresse mail dédiée VSS (style « stop-vss@igensia.fr »)",
             "Pourquoi : signal de réactivité et de simplicité. Comment : créer une adresse, l'afficher sur le site et dans tous les supports apprenants/collaborateurs, communiquer dessus dans le rapport."),
            ("Page dédiée « égalité F/H » avec objectifs chiffrés à 3 ans",
             "Pourquoi : un index qui recule (84→75) doit être suivi d'un plan d'action visible. Comment : 1 page « Notre plan F/H 2025-2027 » avec 3 indicateurs (recrutement, mobilité interne, écarts de salaire) et objectifs annuels."),
        ],
    },
    'BURGUNDY SCHOOL OF BUSINESS': {
        'titre': 'Rapport RSE 2024 (5e édition)',
        'pages': '24 p.',
        'leg_cov': "Couverture sobre, mention « 5e Rapport RSE »",
        'leg_i1':  "Structure 3 piliers : Act for Respect / Sustainability / Empowerment",
        'leg_i2':  "Programme (Re)Connect réfugiés : témoignage + photo",
        'forme_plus': [
            "Structure en 3 piliers (au lieu de 5 axes) : plus mémorable",
            "Photos de programmes spécifiques avec témoignages",
            "Nom de la structure : « Act for Respect / Sustainability / Empowerment » mémorisable",
        ],
        'forme_moins': [
            "Format A4 sobre, manque un peu de souffle visuel",
        ],
        'fond_plus': [
            "20 ans d'engagement RSE (depuis 2005) : argument d'autorité",
            "Module Impact pionnier (1ère grande école française)",
            "Programme (Re)Connect pour les réfugiés (lauréat à plusieurs reprises)",
            "Label DD&RS obtenu juin 2024",
        ],
        'fond_moins': [
            "Pas de bilan carbone dans ce rapport (projet 2024-26)",
            "Budget RSE non publié",
        ],
        'idees': [
            ("Structurer en 3 piliers nommés (verbe d'action) plutôt que 5 axes neutres",
             "Pourquoi : plus mémorable et plus engageant. Comment : trouver 3 verbes-piliers pour IGENSIA (ex : « Former / Inclure / Réduire »). Plus mémorisable que 5 axes du référentiel."),
            ("Mettre en avant 1 ou 2 programmes signature avec témoignages",
             "Pourquoi : ça humanise et différencie. Comment : choisir 2 programmes IGENSIA emblématiques (HOPEN ? handi-accompagnement ?) et leur consacrer 1 double-page chacun avec témoignages."),
        ],
    },
    'EFREI': {
        'titre': 'Rapport de Progrès DD&RS 2026',
        'pages': '23 p.',
        'leg_cov': "Couverture sobre avec accroche « Notre Progrès »",
        'leg_i1':  "Page « bloc de compétences RSE validé par examen »",
        'leg_i2':  "Projets originaux : Mastère Green IT + Harmony + Falcon",
        'forme_plus': [
            "Format compact (23 p.), bien rythmé",
            "Originalités des projets bien mises en avant (chacun avec un visuel)",
        ],
        'forme_moins': [
            "Manque de photos d'apprenants ou collaborateurs",
        ],
        'fond_plus': [
            "Intégration complète DD/RS dans tous les programmes depuis 2024",
            "Bloc de compétences RSE validé par examen (rare et fort)",
            "Charte associations DD obligatoire avec audit carbone",
            "Originalités : Mastère Green IT, projet santé April, Harmony (reconditionnement PC), Falcon (drones EHPAD)",
        ],
        'fond_moins': [
            "Fresque du Climat non mentionnée (déploiement à compléter)",
            "Bilan carbone détaillé non publié dans ce rapport",
        ],
        'idees': [
            ("Imposer une charte DD obligatoire aux associations étudiantes",
             "Pourquoi : crée un effet d'entraînement et matérialise l'engagement étudiant. Comment : signature à la création de l'asso + audit carbone annuel léger des événements organisés."),
            ("Évaluer la compétence RSE de nos apprenants",
             "Pourquoi : ça transforme la RSE de contenu pédagogique en compétence professionnelle. Comment : module évalué (3-5 ECTS) sur les fondamentaux RSE + cas pratique noté."),
        ],
    },
    'EM STRASBOURG BUSINESS SCHOOL': {
        'titre': 'Rapport RSO 2023 (9e édition)',
        'pages': '38 p.',
        'leg_cov': "Couverture austère, mention « 9e Rapport RSO »",
        'leg_i1':  "Budget RSO 92k€ : tableau d'évolution sur 4 ans",
        'leg_i2':  "DU Leadership Méditation Neurosciences avec Matthieu Ricard (photo)",
        'forme_plus': [
            "Tableau d'indicateurs suivis sur 4 années comparatives : trajectoire visible",
            "Mention « 9e Rapport » qui crédibilise l'historique",
        ],
        'forme_moins': [
            "Format austère, peu de mise en récit",
            "Peu de photos qualitatives",
        ],
        'fond_plus': [
            "Budget RSO publié (92k€) avec son évolution sur 4 ans",
            "Label Diversité AFNOR depuis 2012",
            "CforCSR plateforme e-learning obligatoire pour la diplomation",
            "Originalités : Entomovoria FNEGE, B3V cercles de parole, DU avec Matthieu Ricard",
        ],
        'fond_moins': [
            "Bilan carbone propre date de 2010 (trop ancien)",
            "1 seul ETP RSO en baisse",
            "Pas d'objectif GES chiffré",
        ],
        'idees': [
            ("Publier le budget RSE avec son évolution sur 4-5 ans",
             "Pourquoi : c'est l'indicateur de maturité le plus parlant. Comment : tableau simple budget total / ETP dédiés / part du budget global."),
            ("Suivre 10-12 indicateurs RSE sur plusieurs années (vue trajectoire)",
             "Pourquoi : un rapport sans historique sur 3-4 ans = état des lieux, pas démarche. Comment : choisir 10-12 KPIs stables et les republier chaque année (même si évolution faible)."),
        ],
    },
    'PARIS SCHOOL OF BUSINESS\n(Groupe Galileo)': {
        'titre': 'Sustainable Development Progress Report 2023-2024',
        'pages': '19 p.',
        'leg_cov': "Couverture en anglais, sobre",
        'leg_i1':  "Page Sulitest : 50%+ des étudiants français qui passent le test",
        'leg_i2':  "Page « Label STAR » (interne PSB)",
        'forme_plus': [
            "Format compact (19 p.)",
        ],
        'forme_moins': [
            "Rapport en anglais uniquement (ferme une partie du public francophone)",
            "Mise en page peu engageante, tableau ODD générique",
            "Pas de photos qualitatives",
        ],
        'fond_plus': [
            "Sulitest massif : 50%+ des étudiants français qui passent ce test",
            "Fresque du Climat 100% B3+M2 depuis 3 ans (continuité)",
            "Projets humanitaires long-cours : Nosy Komba (10 ans), Guria en Inde",
        ],
        'fond_moins': [
            "Bilan carbone non publié",
            "Budget RSE non publié",
            "RSE portée par une professeure (pas un poste de direction)",
            "Tableau ODD générique sans personnalisation",
        ],
        'idees': [
            ("Adopter le Sulitest comme indicateur de littératie DD étudiante",
             "Pourquoi : c'est un test externe normé, comparable d'une école à l'autre. Comment : intégrer le Sulitest en L2 (passage obligatoire) et publier le score moyen + taux de réussite."),
        ],
    },
    'CESI': {
        'titre': 'Bilan RSE 2023',
        'pages': '57 p.',
        'leg_cov': "Couverture sobre, mention « 4 axes / 18 enjeux »",
        'leg_i1':  "BEGES 2023 ventilé : 2,1 tCO2e/occupant + plan -20%/3 ans",
        'leg_i2':  "Promotion sociale : 75% étudiants non-cadres + 626 étudiants handicap",
        'forme_plus': [
            "Structure très lisible : 4 axes / 18 enjeux / plans 1/3/5 ans",
            "BEGES granulaire (2,1 tCO2e par occupant) bien visualisé",
        ],
        'forme_moins': [
            "57 pages très textuelles, peu de visuels (manque de respirations)",
            "Aucun témoignage",
            "Format austère qui n'engage pas la lecture",
            "Pas de codification claire entre « réalisé » et « annoncé »",
        ],
        'fond_plus': [
            "Périmètre considérable (25 campus, 25 000 étudiants)",
            "BEGES exemplaire : 2,1 tCO2e/occupant + plan transition -20%/3 ans",
            "Promotion sociale forte : 75% étudiants non-cadres, Index F/H 98/100",
            "626 étudiants en situation de handicap accompagnés par 43 référents",
        ],
        'fond_moins': [
            "Beaucoup d'actions encore annoncées plutôt que déployées (charte achats, charte numérique, alimentation, DD&RS)",
            "Forme/Fond très déséquilibrés : contenu solide noyé dans une présentation austère",
        ],
        'idees': [
            ("Publier le BEGES en tCO2e par occupant (apprenant + collaborateur)",
             "Pourquoi : permet la comparaison entre établissements de tailles différentes. CESI : 2,1. IGENSIA : 25 498 / nb occupants = à calculer et publier. Comment : 1 ratio en infographie de tête de page environnement."),
            ("Contre-exemple : un contenu solide ne suffit pas, la forme est décisive",
             "Pourquoi : CESI a noté 4.25/5 en Fond mais 0.75/5 en Forme. Si IGENSIA produit un excellent contenu mais sans soin éditorial, on perd la moitié de l'impact. Investir autant en rédaction que en design."),
        ],
    },
}

# ===== Création prez =====
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W, SLIDE_H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    return shp

def add_text(slide, x, y, w, h, text, *, font='Calibri', size=14, bold=False, color=TEXT, align='left', italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,'right':PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font; run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return box

def add_bullet_list(slide, x, y, w, h, items, *, font='Calibri', size=11, color=TEXT, bullet='•', bullet_color=ROSE):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run(); r1.text = f'{bullet} '
        r1.font.size = Pt(size); r1.font.color.rgb = bullet_color; r1.font.bold = True; r1.font.name = font
        r2 = p.add_run(); r2.text = item
        r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = font
        p.space_after = Pt(3)
    return box

# Compteur de slides — sera mis à jour à la fin
slide_records = []  # liste pour numéroter après

def add_footer(slide, num, total, label='Benchmark RSE — Focus Rapports'):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), VIOLET)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.3), Inches(8), Inches(0.28),
             label, font='Calibri', size=9, color=WHITE, align='left')
    add_text(slide, SLIDE_W - Inches(3), SLIDE_H - Inches(0.3), Inches(2.6), Inches(0.28),
             f'Direction du Développement Durable • mai 2026 • {num}/{total}',
             font='Calibri', size=9, color=WHITE, align='right')

def add_header_band(slide, title, subtitle=None, color=ROSE):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), color)
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), SLIDE_W - Inches(1.3), Inches(0.18),
                                  height=Inches(0.55))
    add_text(slide, Inches(0.4), Inches(0.12), Inches(10.5), Inches(0.45),
             title, font='Calibri', size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.5), Inches(10.5), Inches(0.32),
                 subtitle, font='Calibri', size=11, italic=True, color=WHITE)

def get_orientation(cover_key):
    return meta.get(cover_key, {}).get('orientation', 'portrait')

# Cache des ratios d'images (largeur/hauteur en pixels)
from PIL import Image as PILImage
_ratio_cache = {}
def img_ratio(path):
    p = str(path)
    if p in _ratio_cache: return _ratio_cache[p]
    try:
        with PILImage.open(p) as im:
            w, h = im.size
            r = w / h
    except Exception:
        r = 0.7  # fallback portrait A4
    _ratio_cache[p] = r
    return r

def add_picture_fitted(slide, path, x, y, max_w, max_h, frame_color=None, frame_pad_in=0.04):
    """Insère une image en préservant le ratio, centrée dans la zone (x,y,max_w,max_h).
    Encadre éventuellement avec frame_color (bord)."""
    if not path or not pathlib.Path(path).exists():
        return None
    r_img = img_ratio(path)  # largeur/hauteur
    r_box = max_w / max_h
    if r_img > r_box:
        # image plus large : on fixe la largeur, hauteur calculée
        w = max_w
        h = max_w / r_img
    else:
        # image plus haute : on fixe la hauteur, largeur calculée
        h = max_h
        w = max_h * r_img
    # Centrage dans la zone
    dx = (max_w - w) / 2
    dy = (max_h - h) / 2
    if frame_color is not None:
        pad = Inches(frame_pad_in)
        add_rect(slide, x + dx - pad, y + dy - pad,
                 w + 2*pad, h + 2*pad, frame_color)
    slide.shapes.add_picture(str(path), x + dx, y + dy, width=w, height=h)
    return (x + dx, y + dy, w, h)

# ===== Calcul du nombre total de slides =====
# 1 cover + 1 sommaire + 1 classement + 22*2 + 1 synthèse = 48
TOTAL = 1 + 1 + 1 + len(ranking)*2 + 1
print(f'Total slides prévues : {TOTAL}')

# ============================================================
# SLIDE 1 — Couverture
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, VIOLET)
add_rect(s, 0, 0, Inches(0.4), SLIDE_H, ROSE)
if LOGO_PATH.exists():
    s.shapes.add_picture(str(LOGO_PATH), Inches(1), Inches(0.7), height=Inches(0.8))
add_text(s, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
         'Benchmark', font='Calibri', size=64, bold=True, color=WHITE)
add_text(s, Inches(1), Inches(3.6), Inches(11), Inches(1),
         'Rapports RSE', font='Calibri', size=48, bold=True, color=ROSE)
add_text(s, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
         '22 rapports analysés • Forme et Fond • idées concrètes à transposer chez IGENSIA',
         font='Calibri', size=14, color=WHITE, italic=True)
add_text(s, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
         'Direction du Développement Durable • mai 2026',
         font='Calibri', size=11, color=ROSE_LIGHT)

# ============================================================
# SLIDE 2 — Sommaire
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
add_header_band(s, 'Sommaire', "22 rapports analysés, classés par note finale (Forme + Fond)")

left_items = [(d['rank'], d['short'], d['final']) for nm, d in ranking[:11]]
right_items = [(d['rank'], d['short'], d['final']) for nm, d in ranking[11:]]

def render_sommaire_col(slide, x, items):
    box = slide.shapes.add_textbox(x, Inches(1.2), Inches(6), Inches(5.8))
    tf = box.text_frame; tf.word_wrap = True
    for i, (rank, name, score) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run(); r1.text = f'#{rank:>2d}  '
        r1.font.name = 'Consolas'; r1.font.size = Pt(13); r1.font.bold = True
        is_ig = 'IGENSIA' in name.upper()
        r1.font.color.rgb = ROSE if is_ig else VIOLET
        r2 = p.add_run(); r2.text = name[:35]
        r2.font.name = 'Calibri'; r2.font.size = Pt(13)
        r2.font.color.rgb = ROSE if is_ig else TEXT
        r2.font.bold = is_ig
        if score is not None:
            r3 = p.add_run(); r3.text = f'   {score:.2f}/5'
            r3.font.name = 'Consolas'; r3.font.size = Pt(11); r3.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(4)

render_sommaire_col(s, Inches(0.6), left_items)
render_sommaire_col(s, Inches(6.9), right_items)
add_footer(s, 2, TOTAL)

# ============================================================
# SLIDE 3 — Classement
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
add_header_band(s, 'Classement final',
                f"IGENSIA #{igensia_rank}/{len(ranking)} • Top 3 : OMNES, ESSEC, Galileo")

fig, ax = plt.subplots(figsize=(12, 6), dpi=140)
names = [d['short'] for nm, d in ranking][::-1]
scores = [d['final'] for nm, d in ranking][::-1]
colors = ['#E60F7D' if d['is_igensia'] else '#4A1942' for nm, d in ranking][::-1]
y_pos = range(len(names))
bars = ax.barh(y_pos, scores, color=colors, edgecolor='white', height=0.78)
ax.set_yticks(y_pos)
ax.set_yticklabels(names, fontsize=10)
ax.set_xlim(0, 5.2)
ax.set_xlabel('Note finale /5', fontsize=11, color='#666')
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.spines['left'].set_color('#ccc'); ax.spines['bottom'].set_color('#ccc')
ax.tick_params(colors='#666')
ax.grid(axis='x', linestyle='--', alpha=0.4)
ax.axvline(x=3.5, color='#999', linestyle=':', linewidth=1, alpha=0.5)
for i, (b, v, (nm, d)) in enumerate(zip(bars, scores, ranking[::-1])):
    ax.text(v + 0.06, i, f'{v:.2f}', va='center', fontsize=8.5, fontweight='bold',
            color='#E60F7D' if d['is_igensia'] else '#4A1942')
plt.tight_layout()
chart = pathlib.Path('_pptx_classement.png')
plt.savefig(chart, dpi=140, bbox_inches='tight', facecolor='white')
plt.close()
s.shapes.add_picture(str(chart), Inches(0.5), Inches(1.1), width=Inches(12.3))
add_footer(s, 3, TOTAL)

# ============================================================
# SLIDES 4+ — 2 slides par rapport
# ============================================================
slide_num = 4
for idx, (nm, d) in enumerate(ranking):
    analysis = ANALYSES.get(nm, {})
    cover_key = COVER_KEYS.get(nm)
    orientation = get_orientation(cover_key) if cover_key else 'portrait'
    rank = d['rank']
    short_name = d['short'].replace('★ ','')
    if d['is_igensia']: short_name = '★ ' + short_name
    header_color = ROSE if d['is_igensia'] else VIOLET

    # ----------------------- SLIDE "VISUELS" -----------------------
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    # Header
    add_rect(s, 0, 0, SLIDE_W, Inches(0.85), header_color)
    if LOGO_PATH.exists():
        s.shapes.add_picture(str(LOGO_PATH), SLIDE_W - Inches(1.3), Inches(0.18),
                              height=Inches(0.55))
    add_text(s, Inches(0.4), Inches(0.12), Inches(1.2), Inches(0.6),
             f'#{rank}', font='Calibri', size=30, bold=True, color=WHITE)
    add_text(s, Inches(1.5), Inches(0.13), Inches(8.5), Inches(0.4),
             short_name.upper(), font='Calibri', size=18, bold=True, color=WHITE)
    sub = f'{analysis.get("titre","Rapport")} • {analysis.get("pages","-")}'
    if d['final'] is not None:
        sub += f'  •  Forme {d["avg_f"]:.2f} • Fond {d["avg_b"]:.2f}  •  Note {d["final"]:.2f}/5'
    add_text(s, Inches(1.5), Inches(0.5), Inches(10), Inches(0.32),
             sub, font='Calibri', size=10, color=WHITE, italic=True)
    add_text(s, Inches(0.4), Inches(0.92), Inches(8), Inches(0.4),
             'Le rapport en images', font='Calibri', size=14, bold=True, color=VIOLET)

    # Disposition images : couverture grande à gauche, 2 pages intérieures à droite (empilées)
    cov_path = COVERS_DIR / f'{cover_key}.png' if cover_key else None
    i1_path = COVERS_DIR / f'{cover_key}_inner1.png' if cover_key else None
    i2_path = COVERS_DIR / f'{cover_key}_inner2.png' if cover_key else None

    # Couverture grande à gauche — zone allouée, image centrée préservant ratio
    cov_x = Inches(0.5)
    cov_y = Inches(1.4)
    cov_box_w = Inches(5.0)
    cov_box_h = Inches(5.6)

    if cov_path and cov_path.exists():
        add_picture_fitted(s, cov_path, cov_x, cov_y, cov_box_w, cov_box_h, frame_color=VIOLET)
    else:
        add_rect(s, cov_x, cov_y, cov_box_w, cov_box_h, GREY)
        add_text(s, cov_x, cov_y + cov_box_h/2, cov_box_w, Inches(0.4),
                 '(Couverture non disponible)', font='Calibri', size=11, color=TEXT_LIGHT, align='center', italic=True)
    add_text(s, cov_x, cov_y + cov_box_h + Inches(0.05), cov_box_w, Inches(0.3),
             '◆ ' + (analysis.get('leg_cov','Couverture du rapport')),
             font='Calibri', size=9, italic=True, color=TEXT_LIGHT)

    # 2 pages intérieures à droite (empilées, chacune préserve son ratio)
    right_x = cov_x + cov_box_w + Inches(0.4)
    right_w = SLIDE_W - right_x - Inches(0.5)
    right_y = Inches(1.4)
    right_total_h = Inches(5.6)
    inner_h = (right_total_h - Inches(0.7)) / 2  # 2 zones + légendes

    # Page intérieure 1
    if i1_path and i1_path.exists():
        add_picture_fitted(s, i1_path, right_x, right_y, right_w, inner_h, frame_color=VIOLET)
    add_text(s, right_x, right_y + inner_h + Inches(0.05), right_w, Inches(0.25),
             '◆ ' + (analysis.get('leg_i1','Page intérieure 1')),
             font='Calibri', size=9, italic=True, color=TEXT_LIGHT)

    # Page intérieure 2
    right_y2 = right_y + inner_h + Inches(0.45)
    if i2_path and i2_path.exists():
        add_picture_fitted(s, i2_path, right_x, right_y2, right_w, inner_h, frame_color=VIOLET)
    add_text(s, right_x, right_y2 + inner_h + Inches(0.05), right_w, Inches(0.25),
             '◆ ' + (analysis.get('leg_i2','Page intérieure 2')),
             font='Calibri', size=9, italic=True, color=TEXT_LIGHT)

    add_footer(s, slide_num, TOTAL)
    slide_num += 1

    # ----------------------- SLIDE "ANALYSE" -----------------------
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
    # Header identique
    add_rect(s, 0, 0, SLIDE_W, Inches(0.85), header_color)
    if LOGO_PATH.exists():
        s.shapes.add_picture(str(LOGO_PATH), SLIDE_W - Inches(1.3), Inches(0.18),
                              height=Inches(0.55))
    add_text(s, Inches(0.4), Inches(0.12), Inches(1.2), Inches(0.6),
             f'#{rank}', font='Calibri', size=30, bold=True, color=WHITE)
    add_text(s, Inches(1.5), Inches(0.13), Inches(8.5), Inches(0.4),
             short_name.upper(), font='Calibri', size=18, bold=True, color=WHITE)
    add_text(s, Inches(1.5), Inches(0.5), Inches(10), Inches(0.32),
             'Analyse Forme + Fond • idées concrètes transposables',
             font='Calibri', size=10, color=WHITE, italic=True)

    # 3 colonnes : Forme | Fond | Idées
    col_y = Inches(1.05); col_h = SLIDE_H - Inches(1.4)
    col_w = (SLIDE_W - Inches(1.0) - Inches(0.4)) / 3  # 3 colonnes équilibrées
    gap = Inches(0.2)
    x0 = Inches(0.5)

    forme_plus = analysis.get('forme_plus', [])
    forme_moins = analysis.get('forme_moins', [])
    fond_plus = analysis.get('fond_plus', [])
    fond_moins = analysis.get('fond_moins', [])
    idees = analysis.get('idees', [])

    def render_section(slide, x, y, w, h, title, plus_items, minus_items, accent_color):
        # carte fond
        add_rect(slide, x, y, w, h, WHITE)
        # bandeau titre
        add_rect(slide, x, y, w, Inches(0.45), accent_color)
        add_text(slide, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.3),
                 title, font='Calibri', size=14, bold=True, color=WHITE)
        # +
        inner_y = y + Inches(0.55)
        add_text(slide, x + Inches(0.15), inner_y, w - Inches(0.3), Inches(0.3),
                 '✓ Points forts', font='Calibri', size=11, bold=True, color=GREEN)
        add_bullet_list(slide, x + Inches(0.15), inner_y + Inches(0.3), w - Inches(0.3), Inches(2.2),
                        plus_items, size=10.5, bullet_color=GREEN)
        # –
        inner_y2 = y + Inches(0.55) + Inches(2.7)
        add_text(slide, x + Inches(0.15), inner_y2, w - Inches(0.3), Inches(0.3),
                 '✗ Points faibles', font='Calibri', size=11, bold=True, color=RED)
        add_bullet_list(slide, x + Inches(0.15), inner_y2 + Inches(0.3), w - Inches(0.3), Inches(2.2),
                        minus_items, size=10.5, bullet_color=RED)

    render_section(s, x0, col_y, col_w, col_h, 'Forme', forme_plus, forme_moins, VIOLET)
    render_section(s, x0 + col_w + gap, col_y, col_w, col_h, 'Fond', fond_plus, fond_moins, VIOLET)

    # 3e colonne : Idées à reprendre — format spécial avec titre + "pourquoi/comment"
    x3 = x0 + 2*(col_w + gap)
    add_rect(s, x3, col_y, col_w, col_h, WHITE)
    add_rect(s, x3, col_y, col_w, Inches(0.45), RGBColor(0xD9, 0x9F, 0x00))
    add_text(s, x3 + Inches(0.15), col_y + Inches(0.08), col_w - Inches(0.3), Inches(0.3),
             '★ Idées concrètes pour IGENSIA', font='Calibri', size=14, bold=True, color=WHITE)

    iy = col_y + Inches(0.55)
    for idea in idees[:3]:  # max 3 idées par slide
        title_idea, expl = idea if isinstance(idea, tuple) and len(idea) == 2 else (idea, '')
        # Titre idée
        add_text(s, x3 + Inches(0.15), iy, col_w - Inches(0.3), Inches(0.5),
                 '▸ ' + title_idea, font='Calibri', size=11, bold=True, color=VIOLET)
        # Hauteur titre flexible : 2 lignes max
        iy_after_title = iy + Inches(0.55)
        # Explication
        if expl:
            add_text(s, x3 + Inches(0.3), iy_after_title, col_w - Inches(0.45), Inches(1.2),
                     expl, font='Calibri', size=9.5, color=TEXT, italic=False)
            iy = iy_after_title + Inches(1.25)
        else:
            iy = iy_after_title

    add_footer(s, slide_num, TOTAL)
    slide_num += 1

# ============================================================
# DERNIÈRE SLIDE — Synthèse "Top 5 idées à garder pour IGENSIA"
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BG)
add_header_band(s, 'Cinq idées prioritaires pour notre prochain rapport',
                "Synthèse des bonnes pratiques observées et transposables chez IGENSIA")

idees_top = [
    ('01', "Publier le budget RSE consolidé avec son évolution",
     "Source : OMNES (310k€ depuis 2020), EM Strasbourg (92k€ sur 4 ans). Un chiffre absolu vaut mieux que « 10 ETP »."),
    ('02', "Construire un tableau comparatif annuel des 12-15 KPIs",
     "Source : Audencia (15 indicateurs comparés 2024/2025 en 1 page). Permet de montrer la trajectoire, pas seulement l'état."),
    ('03', "Faire auditer notre bilan carbone par un tiers (AFNOR / EY / Bureau Veritas)",
     "Source : HEC (audit AFNOR), GEM (trajectoire SBTi). Crédibilise notre chiffre 25 498 tCO2e. Coût ~10-15k€ à prévoir."),
    ('04', "Adopter la codification Audencia « Réalisé / Nouveau / En cours » par axe",
     "Source : Audencia + ESSEC. Notre format « C'est réalisé / lancé / prévu » est unique : à systématiser sur tous les axes."),
    ('05', "Intégrer 4-6 témoignages incarnés d'apprenants et collaborateurs",
     "Source : HEC (20 témoignages), EMLYON, BSB. Humanise le rapport et donne des preuves de vie pour chaque axe."),
]
rec_y = Inches(1.2); rec_h = Inches(1.05); rec_gap = Inches(0.12)
for i, (num, title, body) in enumerate(idees_top):
    y = rec_y + i * (rec_h + rec_gap)
    add_rect(s, Inches(0.5), y, Inches(0.9), rec_h, VIOLET)
    add_text(s, Inches(0.5), y + Inches(0.28), Inches(0.9), Inches(0.5),
             num, font='Calibri', size=22, bold=True, color=ROSE, align='center')
    add_rect(s, Inches(1.4), y, SLIDE_W - Inches(1.9), rec_h, WHITE)
    add_text(s, Inches(1.6), y + Inches(0.12), SLIDE_W - Inches(2.1), Inches(0.35),
             title, font='Calibri', size=14, bold=True, color=VIOLET)
    add_text(s, Inches(1.6), y + Inches(0.48), SLIDE_W - Inches(2.1), Inches(0.55),
             body, font='Calibri', size=11, color=TEXT)

add_footer(s, slide_num, TOTAL)

# Sauvegarde
out = pathlib.Path('Benchmark_RSE_IGENSIA_Focus_rapports_v3b.pptx')
prs.save(out)
print(f'OK -- {out.absolute()} ({out.stat().st_size//1024//1024} Mo, {slide_num} slides)')
chart.unlink(missing_ok=True)

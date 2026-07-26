# -*- coding: utf-8 -*-
"""Enrichit le PPTX Benchmark_RSE_DRAFTV6_3.pptx en ajoutant des slides
d'enrichissement à la fin. Rien n'est modifié dans les slides existantes."""
import sys, pathlib
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = pathlib.Path(r'C:\Users\sylva\OneDrive\Bureau\Benchmark_RSE_DRAFTV6_3.pptx')
DST = pathlib.Path(r'C:\Users\sylva\OneDrive\Bureau\Benchmark_RSE_DRAFTV7_enrichi.pptx')

# Palette IGENSIA
VIOLET = RGBColor(0x4A, 0x19, 0x42)
VIOLET_DK = RGBColor(0x26, 0x0D, 0x66)
ROSE = RGBColor(0xE6, 0x0F, 0x7D)
ROSE_LIGHT = RGBColor(0xFD, 0xE3, 0xF1)
BG_LIGHT = RGBColor(0xF8, 0xF6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)
BORDER = RGBColor(0xE0, 0xD8, 0xF0)
GREEN = RGBColor(0x4C, 0xAF, 0x50)

prs = Presentation(str(SRC))
SW = prs.slide_width
SH = prs.slide_height
print(f'Slide size: {SW/914400:.2f} x {SH/914400:.2f} inches')
print(f'Slides existantes : {len(prs.slides)}')

BLANK = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=None):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        if line_w:
            r.line.width = line_w
    r.shadow.inherit = False
    return r

def add_text(slide, x, y, w, h, text, size=12, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri',
             line_spacing=1.2, italic=False):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tx

def add_para(textbox, text, size=10, bold=False, color=TEXT,
             italic=False, line_spacing=1.3, before_pt=2, indent_level=0):
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = line_spacing
    p.space_before = Pt(before_pt)
    if indent_level:
        p.level = indent_level
    r = p.add_run()
    r.text = text
    r.font.name = 'Calibri'
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color

def add_header(slide, axe_label, title):
    """Bandeau de titre violet avec le pilier en surtitre."""
    add_rect(slide, 0, 0, SW, Inches(1.0), fill=VIOLET)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.32),
             axe_label.upper(),
             size=10, bold=True, color=RGBColor(0xE0, 0xD8, 0xF0),
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1)
    add_text(slide, Inches(0.5), Inches(0.45), Inches(12), Inches(0.55),
             title, size=22, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

def add_column(slide, x, y, w, h, header_color, header_text, items, reco=None):
    """Carte colonne avec header + items {name, bullets} + reco optionnelle."""
    # Fond carte
    add_rect(slide, x, y, w, h, fill=BG_LIGHT, line=BORDER, line_w=Emu(9525))
    # Header colonne
    add_rect(slide, x, y, w, Inches(0.45), fill=header_color)
    add_text(slide, x + Inches(0.18), y, w - Inches(0.36), Inches(0.45),
             header_text, size=11, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1)
    # Corps
    body_x = x + Inches(0.18)
    body_y = y + Inches(0.58)
    body_w = w - Inches(0.36)
    body_h = h - Inches(0.7) - (Inches(1.0) if reco else Inches(0.1))
    # Premier item
    first_name, first_bullets = items[0]
    tx = add_text(slide, body_x, body_y, body_w, Inches(0.3),
                  first_name, size=10, bold=True, color=VIOLET, line_spacing=1.25)
    for b in first_bullets:
        add_para(tx, b, size=9, color=TEXT, line_spacing=1.3)
    # Items suivants
    for name, bullets in items[1:]:
        add_para(tx, name, size=10, bold=True, color=VIOLET, before_pt=8, line_spacing=1.25)
        for b in bullets:
            add_para(tx, b, size=9, color=TEXT, line_spacing=1.3)
    # Reco en bas (encadré rose)
    if reco:
        reco_y = y + h - Inches(0.95)
        add_rect(slide, x + Inches(0.1), reco_y, w - Inches(0.2), Inches(0.85),
                 fill=ROSE_LIGHT)
        add_text(slide, x + Inches(0.25), reco_y + Inches(0.06),
                 w - Inches(0.5), Inches(0.25),
                 'RECO IGENSIA', size=8, bold=True, color=ROSE,
                 line_spacing=1)
        add_text(slide, x + Inches(0.25), reco_y + Inches(0.28),
                 w - Inches(0.5), Inches(0.55),
                 reco, size=9, color=TEXT, line_spacing=1.3)

def add_footer(slide, label):
    add_text(slide, Inches(0.5), Inches(7.15), Inches(11), Inches(0.3),
             label, size=8, italic=True, color=TEXT_LIGHT, line_spacing=1)
    # Pavé violet à droite
    add_text(slide, Inches(11.8), Inches(7.15), Inches(1.4), Inches(0.3),
             'Enrichissement V7', size=8, bold=True, color=ROSE,
             align=PP_ALIGN.RIGHT, line_spacing=1)


# ======================================================
# SLIDE 12 — QVT : Santé mentale enrichie
# ======================================================
def slide_qvt_sante_mentale():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Pilier 5 · Qualité de vie & égalité des chances',
               'Santé mentale et innovation pédagogique (à intégrer)')
    # Intro courte
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
             "Notre dispositif For Me (488 entretiens jan-juin 2025) est professionnel et descendant. "
             "Le benchmark révèle 3 leviers complémentaires absents de notre offre — peu coûteux, "
             "fortement reconnus et complémentaires.",
             size=10.5, color=TEXT_LIGHT, line_spacing=1.4)
    # 3 colonnes
    col_w = (SW - Inches(0.8) - Inches(0.4)) / 3
    col_y = Inches(1.75)
    col_h = Inches(5.3)
    for i, col in enumerate([
        {'header': 'Écoute par les pairs', 'color': ROSE,
         'items': [
             ('Nightline (étudiants formés, écoute nocturne)', [
                 '• EDHEC, EMLYON : service gratuit d\'écoute nocturne anonyme',
                 '• Animé par des étudiants formés pour leurs pairs',
                 '• Coût faible (formation), forte adhésion (les jeunes parlent plus à des pairs)',
             ]),
             ('Sentinelles & maillage', [
                 '• EDHEC : réseau d\'étudiants sentinelles formés pour repérer et orienter',
                 '• ESSEC : enseignants, responsables de résidence et bureaux d\'assos formés à détecter les difficultés',
             ]),
         ],
         'reco': 'Former un réseau d\'étudiants-relais par campus, en complément de For Me (détection amont par les pairs).'},
        {'header': 'Soutien 24/7 externalisé', 'color': VIOLET,
         'items': [
             ('Hotline professionnelle 24/7', [
                 '• EDHEC, Pros-Consulte : hotline gratuite tous sujets (stress, anxiété, vie perso)',
                 '• EMLYON, Pôle Santé APICIL : équipe pluridisciplinaire (2 psy, 2 infirmières, médecin) + ligne d\'écoute 7j/7',
             ]),
             ('Lieux dédiés au bien-être', [
                 '• KEDGE, Maison de la Santé et du Bien-Être (oct. 2025) : 1ʳᵉ en France pour une école',
                 '• Regroupe santé physique, mentale, financière et sociale',
             ]),
             ('Crédits académiques', [
                 '• ESSEC, points bien-être crédités académiquement (programme ESSECare, 900 collaborateurs)',
             ]),
         ],
         'reco': 'Compléter For Me par une hotline 24/7 externalisée (hors heures bureau, week-end) et envisager un lieu santé dédié à Nanterre.'},
        {'header': 'Innovation pédagogique', 'color': GREEN,
         'items': [
             ('Réalité virtuelle immersive', [
                 '• EMLYON + Reverto (startup incubée) : Journées de l\'inclusion en VR',
                 '• On vit une situation de handicap, une scène de sexisme au travail',
                 '• Très peu commun, mémorable — argument COMEX et rapport',
             ]),
             ('Outils pédagogiques distinctifs', [
                 '• EDHEC : jeu de cartes de sensibilisation VSS (créé par la Chaire D&I)',
                 '• EDHEC : journée VSS obligatoire à la rentrée pour tous les programmes',
             ]),
         ],
         'reco': 'Tester une journée de l\'inclusion en VR dans HOPEN — 1 expérience marquante par promotion.'},
    ]):
        x = Inches(0.4) + (col_w + Inches(0.2)) * i
        add_column(s, x, col_y, col_w, col_h, col['color'], col['header'],
                   col['items'], reco=col['reco'])
    add_footer(s, 'À intégrer en Slide 10 — Pilier QVT & égalité')


# ======================================================
# SLIDE 13 — Bourses ciblées distinctives
# ======================================================
def slide_bourses_ciblees():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Pilier 5 · Fonds de dotation',
               'Bourses ciblées distinctives (à ajouter à la liste « financer quoi »)')
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
             "La slide actuelle cite bourses sociales, fonds d\'urgence, projets impact, SAFE remboursable. "
             "Le benchmark révèle 3 mécanismes complémentaires distinctifs, alignés avec notre ADN d\'inclusion.",
             size=10.5, color=TEXT_LIGHT, line_spacing=1.4)
    col_w = (SW - Inches(0.8) - Inches(0.4)) / 3
    col_y = Inches(1.75)
    col_h = Inches(5.3)
    for i, col in enumerate([
        {'header': 'Bourses réfugiés', 'color': ROSE,
         'items': [
             ('4 écoles ont un dispositif dédié', [
                 '• HEC, Imagine Fellows : bourses complètes pour étudiants réfugiés et issus de zones de conflit',
                 '• BSB, (Re)Connect : programme de réintégration des réfugiés, déjà primé deux fois',
                 '• GEM, Refugee Grant Program',
                 '• PSB, Each One : réinsertion de réfugiés, piloté avec une association',
             ]),
             ('Brest BS, Projet Tremplin', [
                 '• Avec Fondation Mazars et Each One',
                 '• Formation 3 mois pour réfugiés (FLE 12h/sem)',
             ]),
         ],
         'reco': 'Lancer une bourse réfugiés (volumes modestes, fort signal). Aligné avec notre partenariat ANLCI (illettrisme/illectronisme).'},
        {'header': 'Prêt étudiant co-garanti', 'color': VIOLET,
         'items': [
             ('KEDGE, KAP Equal Opportunities', [
                 '• Prêt co-garanti : école 50 % + banque partenaire 50 %',
                 '• Lève la barrière du financement initial',
                 '• Coût direct faible pour l\'école (garantie, pas dépense)',
                 '• Mécanisme très peu répandu = différenciant fort',
             ]),
             ('AD Education, Fonds SAFE', [
                 '• 750 K€ sous forme d\'avance remboursable une fois en emploi',
                 '• Logique similaire : pas une bourse, un soutien remboursé',
             ]),
         ],
         'reco': 'Étudier un prêt d\'honneur co-garanti avec une banque partenaire, fléché sur les apprenants en difficulté.'},
        {'header': 'Bourse mobilité bas carbone', 'color': GREEN,
         'items': [
             ('ESSEC, budget mobilité durable', [
                 '• Finance les déplacements bas carbone (train vs avion)',
                 '• 45 déplacements financés en 2023/24',
             ]),
             ('EFREI, 61 bourses mobilité bas carbone', [
                 '• Pour soutenir le train ou le bus vers les destinations européennes',
             ]),
             ('EMLYON, chèque mobilité douce', [
                 '• Pour étudiants en mobilité internationale qui choisissent le transport bas carbone',
             ]),
             ('TSM, aide mobilité verte ; Sciences Po, vols courts interdits', []),
         ],
         'reco': 'Créer une « bourse train » pour les stages et échanges, qui croise notre charte mobilité et notre engagement carbone (-40 % d\'ici 2030).'},
    ]):
        x = Inches(0.4) + (col_w + Inches(0.2)) * i
        add_column(s, x, col_y, col_w, col_h, col['color'], col['header'],
                   col['items'], reco=col['reco'])
    add_footer(s, 'À intégrer en Slide 10 — bloc « Fonds de dotation : pour financer quoi exactement »')


# ======================================================
# SLIDE 14 — Gouvernance enrichie
# ======================================================
def slide_gouvernance_enrich():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Pilier 4 · Gouvernance responsable',
               'Vérification externe, exemplarité interne et rémunération indexée')
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
             "La slide gouvernance couvre l\'articulation, la compliance et l\'association des étudiants. "
             "3 angles peu commun à ajouter pour pousser la crédibilité externe et l\'exemplarité interne.",
             size=10.5, color=TEXT_LIGHT, line_spacing=1.4)
    col_w = (SW - Inches(0.8) - Inches(0.4)) / 3
    col_y = Inches(1.75)
    col_h = Inches(5.3)
    for i, col in enumerate([
        {'header': 'Vérification externe (OTI)', 'color': ROSE,
         'items': [
             ('GEM, audit KPMG', [
                 '• Indicateurs RSE audités par KPMG (organisme tiers indépendant)',
                 '• Argument de crédibilité ultime pour un rapport',
             ]),
             ('EMLYON, accréditation COFRAC', [
                 '• Comité de Mission validé par un OTI accrédité COFRAC (rare et exigeant)',
             ]),
             ('EMLYON, label HRS4R', [
                 '• Human Resource Strategy for Researchers (Commission Européenne, fév. 2025)',
                 '• Engagement européen sur la gestion des chercheurs',
             ]),
         ],
         'reco': 'Faire vérifier dès le prochain rapport au moins 3 à 5 indicateurs RSE phares par un organisme tiers. À séquencer tôt (le processus est long).'},
        {'header': 'Mixité de la gouvernance publiée', 'color': VIOLET,
         'items': [
             ('EMLYON, exemplarité chiffrée', [
                 '• COMEX : 40 % de femmes (4/10)',
                 '• CoDir : 47 % de femmes (+18 % depuis 2021)',
                 '• Conseil de Surveillance : 41 %',
                 '• Corps professoral : 37 %',
                 '• Effectifs France : 66 % de femmes',
             ]),
             ('EDHEC, % de présidentes d\'associations', [
                 '• 47 % de présidentes — KPI publié, suivi dans le temps',
             ]),
             ('Index F/H publié en détail', [
                 '• EMLYON : écart rémunération 35/40, augmentations 20/20, promotions 15/15, retour congé maternité 15/15',
             ]),
         ],
         'reco': 'Publier la mixité par instance de gouvernance et le détail de notre index égalité (notre 76/100 mérite d\'être explicité, pas masqué).'},
        {'header': 'Rémunération dirigeante indexée ESG', 'color': GREEN,
         'items': [
             ('AD Education', [
                 '• Les primes du DG sont indexées sur les critères ESG',
                 '• Signal de gouvernance le plus fort du benchmark',
                 '• Très peu pratiqué (1 entité)',
             ]),
             ('Pratique exigeante mais distinctive', [
                 '• Engage la direction sur des résultats RSE mesurables',
                 '• À l\'inverse du déclaratif, force l\'action',
             ]),
         ],
         'reco': 'À explorer pour la trajectoire 2027-2028 : associer une part de la rémunération variable du COMEX à l\'atteinte d\'objectifs RSE chiffrés (déjà alignés sur nos engagements LUCIE).'},
    ]):
        x = Inches(0.4) + (col_w + Inches(0.2)) * i
        add_column(s, x, col_y, col_w, col_h, col['color'], col['header'],
                   col['items'], reco=col['reco'])
    add_footer(s, 'À intégrer en Slide 9 — Pilier Gouvernance')


# ======================================================
# SLIDE 15 — Ancrage territorial enrichi
# ======================================================
def slide_ancrage_enrich():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Pilier 1 · Ancrage territorial / Utilité sociétale',
               'Mécénat de compétences, open badges et labels associations')
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
             "La slide couvre projets citoyens, compétence étudiants/territoire, achat local. "
             "3 leviers absents : l\'engagement des collaborateurs sur leur temps, et la valorisation/labellisation de l\'engagement (côté étudiants).",
             size=10.5, color=TEXT_LIGHT, line_spacing=1.4)
    col_w = (SW - Inches(0.8) - Inches(0.4)) / 3
    col_y = Inches(1.75)
    col_h = Inches(5.3)
    for i, col in enumerate([
        {'header': 'Mécénat de compétences', 'color': ROSE,
         'items': [
             ('Bénévolat sur le temps de travail', [
                 '• HEC Stand Up : programme de 10 semaines pour l\'entrepreneuriat social, accessible au personnel',
                 '• EMLYON : les collaborateurs peuvent consacrer du temps à des associations partenaires',
             ]),
             ('Logique « employeur engagé »', [
                 '• Coût faible (heures, pas budget direct)',
                 '• Valorisation interne et signal d\'engagement',
                 '• Complète notre charte achats responsables côté ressources humaines',
             ]),
         ],
         'reco': 'Offrir aux collaborateurs un quota d\'heures/an pour faire du bénévolat sur le temps de travail (associations partenaires : Restos, Mozaïk, ANLCI, Missions Locales).'},
        {'header': 'Open badges (valoriser l\'engagement)', 'color': VIOLET,
         'items': [
             ('GEM, Olympiades de la Sustainability', [
                 '• 15 jours d\'activités RSE au choix avec délivrance d\'open badges',
                 '• Certificats numériques de soft skills (CV, LinkedIn)',
             ]),
             ('Reconnaissance numérique de compétences', [
                 '• Standard ouvert IMS, transférable',
                 '• Valorise les soft skills RSE acquises hors cursus',
                 '• Sujet émergent, peu d\'écoles le font (= différenciant)',
             ]),
         ],
         'reco': 'Délivrer des open badges via la plateforme MyIGENSIA pour valider l\'engagement RSE des apprenants (HOPEN, projets citoyens, missions associatives).'},
        {'header': 'Labels RSE des associations étudiantes', 'color': GREEN,
         'items': [
             ('EDHEC, label Citizen Associations', [
                 '• Cahier des charges 3 critères (éthique/gouvernance, social, environnemental)',
                 '• Objectif : 100 % des associations labellisées d\'ici 2028',
             ]),
             ('EMLYON, charte RSE des associations', [
                 '• Signée par 42 associations',
                 '• Bilan carbone des 5 événements majeurs, référent VSS, label Ecollab',
                 '• 64 % des assos ont obtenu Ecollab',
             ]),
             ('GEM, représentants RSE par association', [
                 '• 22 associations, ~700 événements/an',
                 '• Identifié comme point fort par l\'audit DD&RS',
             ]),
         ],
         'reco': 'Lancer une charte (ou label) RSE des associations IGENSIA dans le cadre de notre structuration assos prévue déc. 2026 : bilan, référent VSS, formation des bureaux.'},
    ]):
        x = Inches(0.4) + (col_w + Inches(0.2)) * i
        add_column(s, x, col_y, col_w, col_h, col['color'], col['header'],
                   col['items'], reco=col['reco'])
    add_footer(s, 'À intégrer en Slide 4 — Pilier Ancrage territorial')


# ======================================================
# SLIDE 16 — Environnement, compléments distinctifs
# ======================================================
def slide_env_enrich():
    s = prs.slides.add_slide(BLANK)
    add_header(s, 'Pilier 3 · Réduire l\'impact environnemental',
               'Signal-prix carbone et reconditionnement IT chiffré')
    add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.45),
             "La slide environnement traite alimentation, numérique, mobilité, biodiversité et déchets. "
             "2 pratiques distinctives à pousser : la tarification carbone des plats (nudge), "
             "et la quantification du CO2 évité par le reconditionnement IT.",
             size=10.5, color=TEXT_LIGHT, line_spacing=1.4)
    col_w = (SW - Inches(0.8) - Inches(0.2)) / 2
    col_y = Inches(1.75)
    col_h = Inches(5.3)
    for i, col in enumerate([
        {'header': 'Tarification carbone des plats (signal-prix)', 'color': GREEN,
         'items': [
             ('HEC, prix carbone sur les plats', [
                 '• Le plat le plus émetteur coûte plus cher',
                 '• Le plat végétarien est subventionné',
                 '• Effet documenté : 2 jours sans viande rouge/semaine = 278 tCO2 évitées',
                 '• Nudge comportemental peu commun et très narratif',
             ]),
             ('ESSEC, repas végétarien à prix réduit', [
                 '• Subvention spécifique introduite en 2024',
                 '• Logique de signal-prix complémentaire',
             ]),
             ('Pourquoi c\'est intéressant', [
                 '• Pas une interdiction, un choix libre orienté économiquement',
                 '• Effet mesurable, communicable au COMEX',
                 '• Demande de cadrer le contrat prestataire restauration',
             ]),
         ],
         'reco': 'Inclure dans le cahier des charges du prestataire restauration un signal-prix carbone (différentiel de tarif), avec une cible mesurable de tCO2 évitées par an.'},
        {'header': 'Reconditionnement IT chiffré CO2', 'color': VIOLET,
         'items': [
             ('HEC, quantification fine', [
                 '• 5,9 tCO2 évitées par le reconditionnement IT',
                 '• 2,2 tCO2 évitées via le remplacement de batteries',
                 '• Recettes de revente reversées au Sustainability Fund',
             ]),
             ('Boucle vertueuse triple', [
                 '• Allongement durée de vie matériel à 5 ans (EDHEC, EMLYON)',
                 '• Reconditionnement via une structure ESS (EDHEC : ATF Gaia ; Galileo : ZACK / Ateliers Sans Frontières)',
                 '• Insertion par l\'emploi des personnes en situation de handicap (CESI : HANDIPRINT)',
             ]),
             ('IÉSEG, 80 % reconditionné via l\'ESS', [
                 '• KPI publié, ancrage clair dans l\'économie sociale et solidaire',
             ]),
         ],
         'reco': 'Étendre notre Digital Clean Up (1,83 tCO2e évitée) à une politique IT complète : 5 ans + reconditionnement ESS + quantification systématique des tCO2 évitées.'},
    ]):
        x = Inches(0.4) + (col_w + Inches(0.2)) * i
        add_column(s, x, col_y, col_w, col_h, col['color'], col['header'],
                   col['items'], reco=col['reco'])
    add_footer(s, 'À intégrer en Slides 7-8 — Pilier Environnement')


# Génération
slide_qvt_sante_mentale()
slide_bourses_ciblees()
slide_gouvernance_enrich()
slide_ancrage_enrich()
slide_env_enrich()

prs.save(str(DST))
print(f'\nNouveau total : {len(prs.slides)} slides')
print(f'✓ Sauvegardé : {DST}')

# -*- coding: utf-8 -*-
"""V4 — Generator du PPTX Benchmark des Rapports RSE.

Lit la source de vérité actuelle : presentation/data.js
3 slides par fiche (visuels / analyse / idées) + IGENSIA Ref + intro + conclu.
"""
import sys, re, pathlib, json5
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== Charte IGENSIA =====
VIOLET     = RGBColor(0x4A, 0x19, 0x42)
VIOLET_DK  = RGBColor(0x26, 0x0D, 0x66)
ROSE       = RGBColor(0xE6, 0x0F, 0x7D)
ROSE_LIGHT = RGBColor(0xFD, 0xE3, 0xF1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BG         = RGBColor(0xFA, 0xFA, 0xFA)
BG_LIGHT   = RGBColor(0xF8, 0xF6, 0xFF)
BORDER     = RGBColor(0xE0, 0xD8, 0xF0)
TEXT       = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_LIGHT = RGBColor(0x66, 0x66, 0x66)
GREEN      = RGBColor(0x4C, 0xAF, 0x50)
GREEN_BG   = RGBColor(0xE8, 0xF5, 0xE9)
RED        = RGBColor(0xE5, 0x39, 0x35)
RED_BG     = RGBColor(0xFF, 0xEB, 0xEE)
YELLOW     = RGBColor(0xFF, 0xC1, 0x07)
YELLOW_BG  = RGBColor(0xFF, 0xF8, 0xE1)
GREY       = RGBColor(0xEC, 0xEC, 0xEC)

# ===== Load data =====
raw = pathlib.Path('presentation/data.js').read_text(encoding='utf-8')
m = re.search(r'const\s+PRESENTATION_DATA\s*=\s*(\{.*?\});\s*$', raw, re.DOTALL)
if not m:
    sys.exit('Erreur : impossible de parser presentation/data.js')
DATA = json5.loads(m.group(1))
RAPPORTS = DATA['rapports']

IMG_DIR = pathlib.Path('presentation/images')

# ===== Presentation setup (16:9) =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ===== Helpers =====
def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        if line_w is not None:
            rect.line.width = line_w
    rect.shadow.inherit = False
    return rect

def add_text(slide, x, y, w, h, text, size=12, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Calibri',
             line_spacing=1.15):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tx

def add_bullets(slide, x, y, w, h, items, size=11, color=TEXT, font='Calibri',
                bullet_color=None, line_spacing=1.25, max_lines=None):
    """Affiche items comme bullets avec • en couleur d'accent."""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top  = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    items_to_render = items if max_lines is None else items[:max_lines]
    for i, item in enumerate(items_to_render):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        r1 = p.add_run()
        r1.text = '• '
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color or color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    if max_lines is not None and len(items) > max_lines:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = f'… (+ {len(items) - max_lines} autres)'
        r.font.name = font
        r.font.size = Pt(size - 1)
        r.font.italic = True
        r.font.color.rgb = TEXT_LIGHT
    return tx

def add_footer(slide, num, total, label=None):
    """Pied de page : numéro + (optionnel) titre fiche."""
    footer = label or 'Benchmark RSE IGENSIA · Rapports'
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.3),
             footer, size=9, color=TEXT_LIGHT, font='Calibri')
    add_text(slide, Inches(12), Inches(7.05), Inches(0.93), Inches(0.3),
             f'{num} / {total}', size=9, color=TEXT_LIGHT,
             align=PP_ALIGN.RIGHT, font='Calibri')

def add_title_bar(slide, title, subtitle=None, rank=None, score=None):
    """Bandeau de titre coloré en haut de slide."""
    add_rect(slide, 0, 0, SW, Inches(0.95), fill=VIOLET)
    if rank is not None:
        # Pastille rang
        add_rect(slide, Inches(0.4), Inches(0.2), Inches(0.85), Inches(0.55),
                 fill=ROSE)
        add_text(slide, Inches(0.4), Inches(0.2), Inches(0.85), Inches(0.55),
                 f'#{rank}', size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1)
        title_x = Inches(1.4)
    else:
        title_x = Inches(0.4)
    title_w = SW - title_x - (Inches(2.2) if score is not None else Inches(0.4))
    add_text(slide, title_x, Inches(0.12), title_w, Inches(0.5),
             title, size=22, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    if subtitle:
        add_text(slide, title_x, Inches(0.55), title_w, Inches(0.35),
                 subtitle, size=11, color=RGBColor(0xE0, 0xD8, 0xF0),
                 anchor=MSO_ANCHOR.TOP, line_spacing=1)
    if score is not None:
        # Score block à droite
        add_text(slide, SW - Inches(2.1), Inches(0.15), Inches(1.7), Inches(0.65),
                 f'{score:.1f}', size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1)
        add_text(slide, SW - Inches(2.1), Inches(0.5), Inches(1.7), Inches(0.4),
                 '/ 5', size=12, color=RGBColor(0xE0, 0xD8, 0xF0),
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1)

def truncate(s, n=180):
    return s if len(s) <= n else s[:n - 1].rstrip() + '…'


# ====================
# COVER SLIDE
# ====================
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # Fond violet
    add_rect(s, 0, 0, SW, SH, fill=VIOLET)
    # Bande rose accent
    add_rect(s, Inches(0.6), Inches(2.6), Inches(0.15), Inches(2),
             fill=ROSE)
    # Titre
    add_text(s, Inches(0.95), Inches(2.5), Inches(11), Inches(1.2),
             'Benchmark des Rapports RSE',
             size=44, bold=True, color=WHITE, line_spacing=1.1)
    # Sous-titre
    add_text(s, Inches(0.95), Inches(3.7), Inches(11), Inches(0.8),
             'Analyse comparée de 21 rapports concurrents + référence IGENSIA',
             size=20, color=RGBColor(0xE0, 0xD8, 0xF0), line_spacing=1.2)
    # Footer date
    add_text(s, Inches(0.95), Inches(4.6), Inches(11), Inches(0.5),
             'IGENSIA Education · mai 2026',
             size=14, color=ROSE_LIGHT, line_spacing=1.2)
    # Mention
    add_text(s, Inches(0.95), Inches(6.8), Inches(11), Inches(0.4),
             'Travail exploratoire interne',
             size=10, color=RGBColor(0xC0, 0xB0, 0xD8))


# ====================
# LECTURE / MÉTHODO
# ====================
def slide_methodo(num, total):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, 'Comment lire ce benchmark', 'Méthodologie et clés de lecture')
    y = Inches(1.4)
    add_text(s, Inches(0.5), y, Inches(12.3), Inches(0.5),
             'Sources et corpus', size=16, bold=True, color=VIOLET)
    add_bullets(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(1.5),
                ['21 rapports RSE / DD&RS / Impact / Société à Mission publiés par des écoles et groupes concurrents, croisés avec la référence IGENSIA',
                 'Chaque rapport est analysé en lecture intégrale : couverture + pages intérieures distinctives, structure, KPIs, témoignages, signature',
                 'Évaluation indépendante par 2 lecteurs (Forme + Fond) sur une échelle de 0 à 5'],
                size=12)
    add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
             'Structure de chaque fiche', size=16, bold=True, color=VIOLET)
    add_bullets(s, Inches(0.6), Inches(4.15), Inches(12.2), Inches(1.5),
                ['Slide 1 : couverture HD + pages intérieures distinctives avec leurs légendes',
                 'Slide 2 : analyse Forme (visuel, structure, narration) et Fond (engagements, KPIs, gouvernance)',
                 'Slide 3 : idées éditoriales transposables pour le prochain rapport IGENSIA'],
                size=12)
    add_text(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(0.5),
             'Lecture des codes couleur', size=16, bold=True, color=VIOLET)
    # Légende couleurs
    legend_y = Inches(6.3)
    add_rect(s, Inches(0.6), legend_y, Inches(0.3), Inches(0.3),
             fill=GREEN_BG, line=GREEN, line_w=Emu(9525))
    add_text(s, Inches(1.0), legend_y - Inches(0.02), Inches(4), Inches(0.35),
             'Points forts (Forme / Fond)', size=11, color=TEXT,
             anchor=MSO_ANCHOR.TOP)
    add_rect(s, Inches(4.6), legend_y, Inches(0.3), Inches(0.3),
             fill=RED_BG, line=RED, line_w=Emu(9525))
    add_text(s, Inches(5.0), legend_y - Inches(0.02), Inches(4), Inches(0.35),
             'Points faibles (Forme / Fond)', size=11, color=TEXT)
    add_rect(s, Inches(8.6), legend_y, Inches(0.3), Inches(0.3),
             fill=ROSE_LIGHT, line=ROSE, line_w=Emu(9525))
    add_text(s, Inches(9.0), legend_y - Inches(0.02), Inches(4), Inches(0.35),
             'Idées à reprendre pour IGENSIA', size=11, color=TEXT)
    add_footer(s, num, total, 'Méthodologie')


# ====================
# IGENSIA REFERENCE
# ====================
def slide_igensia_intro(num, total):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, 'Notre rapport — Rapport RSE Groupe IGENSIA 2024-2025',
                  '48 pages · 1er rapport RSE publié par le Groupe')
    # Cover image (left)
    cov = IMG_DIR / 'igensia.jpg'
    if cov.exists():
        s.shapes.add_picture(str(cov), Inches(0.5), Inches(1.2),
                             height=Inches(5.5))
    # Inner images (right grid)
    inners = ['igensia_inner1.jpg', 'igensia_inner2.jpg']
    legs = [
        'Timeline du parcours RSE 2023-2025',
        'ODD et labellisation LUCIE 26000',
    ]
    inner_x = Inches(4.8)
    inner_y0 = Inches(1.2)
    img_h = Inches(2.6)
    img_w = Inches(4.05)
    for i, (img, leg) in enumerate(zip(inners, legs)):
        p = IMG_DIR / img
        if p.exists():
            y = inner_y0 + (img_h + Inches(0.4)) * i
            s.shapes.add_picture(str(p), inner_x, y, height=img_h)
            add_text(s, inner_x + img_w + Inches(0.2),
                     y + Inches(0.2), Inches(4.3), img_h,
                     leg, size=10, color=TEXT, line_spacing=1.3)
    add_footer(s, num, total, 'IGENSIA Reference')


def slide_igensia_forces(num, total):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, 'IGENSIA — Ce qui marche', 'Forces identifiées sur notre rapport actuel')
    # Forme
    add_rect(s, Inches(0.4), Inches(1.2), Inches(6.3), Inches(5.5),
             fill=GREEN_BG, line=GREEN, line_w=Emu(9525))
    add_text(s, Inches(0.65), Inches(1.35), Inches(5.8), Inches(0.4),
             '✓ Forme', size=18, bold=True, color=GREEN)
    add_bullets(s, Inches(0.65), Inches(1.85), Inches(5.8), Inches(4.7),
                [
                    'Codification REALISE / LANCE / PREVU : structurante, rare dans le benchmark',
                    'Photos qualitatives en pleine page, schémas et illustrations soignés',
                    '8 témoignages incarnés (étudiant, collaborateurs, directeurs, partenaires)',
                    'Structure mémorisable : 4 grandes parties Apprenants / Collaborateurs / Campus / Partenaires',
                ],
                size=12, color=TEXT, bullet_color=GREEN)
    # Fond
    add_rect(s, Inches(6.95), Inches(1.2), Inches(6.0), Inches(5.5),
             fill=GREEN_BG, line=GREEN, line_w=Emu(9525))
    add_text(s, Inches(7.2), Inches(1.35), Inches(5.5), Inches(0.4),
             '✓ Fond', size=18, bold=True, color=GREEN)
    add_bullets(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(4.7),
                [
                    'Bilan Carbone détaillé',
                    'Double signature gouvernance : DG + Directrice DD',
                    'Gouvernance RSE structurée : 3 niveaux + 17 pilotes nommés',
                    'Label LUCIE 26000 + ancrage ODD : tiers externes crédibilisants',
                    'Section « C\'EST PRÉVU » avec objectifs datés (rentrée 2027 : 90 % apprenants 1A/2A/3A formés)',
                ],
                size=12, color=TEXT, bullet_color=GREEN)
    add_footer(s, num, total, 'IGENSIA Reference · Forces')


def slide_igensia_reco(num, total):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, 'IGENSIA — Ce qu\'on peut améliorer', 'Recommandations issues du benchmark')
    # Forme
    add_rect(s, Inches(0.4), Inches(1.2), Inches(6.3), Inches(5.7),
             fill=RED_BG, line=RED, line_w=Emu(9525))
    add_text(s, Inches(0.65), Inches(1.35), Inches(5.8), Inches(0.4),
             '✗ Forme', size=18, bold=True, color=RED)
    add_bullets(s, Inches(0.65), Inches(1.85), Inches(5.8), Inches(4.9),
                [
                    'Label LUCIE 26000 pas suffisamment visible sur la page de garde',
                    'Pas de schéma directeur RSE 5 axes visible en début de rapport',
                    'Pas de page « Chiffres-clés » synthétique en ouverture',
                    'Pas de frise historique de l\'engagement RSE du Groupe',
                    'Codification REALISE/LANCE/PREVU trop binaire sur les objectifs datés',
                ],
                size=11, color=TEXT, bullet_color=RED, line_spacing=1.3)
    # Fond
    add_rect(s, Inches(6.95), Inches(1.2), Inches(6.0), Inches(5.7),
             fill=RED_BG, line=RED, line_w=Emu(9525))
    add_text(s, Inches(7.2), Inches(1.35), Inches(5.5), Inches(0.4),
             '✗ Fond', size=18, bold=True, color=RED)
    add_bullets(s, Inches(7.2), Inches(1.85), Inches(5.5), Inches(4.9),
                [
                    'Pas de budget RSE publié en chiffre absolu',
                    'Pas de tableau récapitulatif d\'indicateurs N / N-1',
                    'ODD pas explicitement connectés à nos actions',
                    'LUCIE 26000 sous-exploitée comme épine dorsale du rapport',
                    'Impact carbone évité non chiffré par action',
                    'Pas d\'auto-critique structurée des défis à relever',
                    'Dimension Groupe sous-exploitée',
                ],
                size=11, color=TEXT, bullet_color=RED, line_spacing=1.3)
    add_footer(s, num, total, 'IGENSIA Reference · Recommandations')


# ====================
# RAPPORT FICHE — 3 slides
# ====================
def slide_rapport_visuels(r, num, total):
    """Slide A : header + cover + inner images avec légendes."""
    s = prs.slides.add_slide(BLANK)
    sub = f"{r['titre']} · {r['pages']}"
    add_title_bar(s, r['name'], sub, rank=r['rank'], score=r['score'])
    # Sous-bandeau scores
    add_text(s, Inches(0.5), Inches(1.1), Inches(7), Inches(0.4),
             f"Forme {r['forme']:.2f} · Fond {r['fond']:.2f}",
             size=11, bold=True, color=ROSE)
    # Cover image (left, large)
    cov = IMG_DIR / f"{r['key']}.jpg"
    cover_landscape = r.get('cover_landscape', False)
    if cov.exists():
        if cover_landscape:
            # Landscape : pleine largeur en haut
            s.shapes.add_picture(str(cov), Inches(0.5), Inches(1.55),
                                 width=Inches(6.2))
            cover_h = Inches(3.5)
        else:
            # Portrait : à gauche, hauteur 5 inches
            s.shapes.add_picture(str(cov), Inches(0.5), Inches(1.55),
                                 height=Inches(5.0))
            cover_h = Inches(5.0)
    # Légende cover (sous l'image)
    add_text(s, Inches(0.5), Inches(6.62), Inches(6.2), Inches(0.4),
             truncate(r.get('leg_cov', ''), 180),
             size=8, color=TEXT_LIGHT, line_spacing=1.25)
    # Inner images (right grid) — 4 max, 2x2
    inners = []
    for i in range(1, 6):
        leg = r.get(f'leg_i{i}')
        if leg:
            img = IMG_DIR / f"{r['key']}_inner{i}.jpg"
            if img.exists():
                inners.append((img, leg))
    inners = inners[:4]
    grid_x = Inches(7.0)
    grid_y = Inches(1.55)
    cell_w = Inches(3.05)
    cell_h = Inches(2.55)
    gap = Inches(0.15)
    for i, (img, leg) in enumerate(inners):
        row = i // 2
        col = i % 2
        x = grid_x + (cell_w + gap) * col
        y = grid_y + (cell_h + gap) * row
        # Image (top 60%)
        img_h = Inches(1.55)
        s.shapes.add_picture(str(img), x, y, width=cell_w, height=img_h)
        # Légende dessous
        add_text(s, x, y + img_h + Inches(0.05), cell_w, cell_h - img_h - Inches(0.05),
                 truncate(leg, 130), size=7, color=TEXT_LIGHT,
                 line_spacing=1.2)
    add_footer(s, num, total, f"{r['name']} · Visuels")


def slide_rapport_analyse(r, num, total):
    """Slide B : Forme + Fond (+/−)."""
    s = prs.slides.add_slide(BLANK)
    sub = f"{r['titre']} · {r['pages']}"
    add_title_bar(s, f"{r['name']} — Analyse", sub, rank=r['rank'], score=r['score'])

    # Limites pour ne pas saturer
    MAX_BULLETS = 6

    # Forme — colonne gauche
    col_x = Inches(0.4)
    col_w = Inches(6.3)
    add_rect(s, col_x, Inches(1.25), col_w, Inches(5.65),
             fill=BG_LIGHT, line=BORDER, line_w=Emu(9525))
    add_text(s, col_x + Inches(0.25), Inches(1.4), col_w - Inches(0.5),
             Inches(0.35),
             'FORME', size=11, bold=True, color=VIOLET)
    # Forme +
    add_text(s, col_x + Inches(0.25), Inches(1.8), col_w - Inches(0.5),
             Inches(0.35),
             '✓ Points forts', size=12, bold=True, color=GREEN)
    add_bullets(s, col_x + Inches(0.25), Inches(2.2), col_w - Inches(0.5),
                Inches(2.0), r.get('forme_plus', []),
                size=10, color=TEXT, bullet_color=GREEN,
                line_spacing=1.25, max_lines=MAX_BULLETS)
    # Forme −
    add_text(s, col_x + Inches(0.25), Inches(4.3), col_w - Inches(0.5),
             Inches(0.35),
             '✗ Points faibles', size=12, bold=True, color=RED)
    add_bullets(s, col_x + Inches(0.25), Inches(4.7), col_w - Inches(0.5),
                Inches(2.1), r.get('forme_moins', []),
                size=10, color=TEXT, bullet_color=RED,
                line_spacing=1.25, max_lines=MAX_BULLETS)

    # Fond — colonne droite
    col_x = Inches(6.95)
    col_w = Inches(6.0)
    add_rect(s, col_x, Inches(1.25), col_w, Inches(5.65),
             fill=BG_LIGHT, line=BORDER, line_w=Emu(9525))
    add_text(s, col_x + Inches(0.25), Inches(1.4), col_w - Inches(0.5),
             Inches(0.35),
             'FOND', size=11, bold=True, color=VIOLET)
    # Fond +
    add_text(s, col_x + Inches(0.25), Inches(1.8), col_w - Inches(0.5),
             Inches(0.35),
             '✓ Points forts', size=12, bold=True, color=GREEN)
    add_bullets(s, col_x + Inches(0.25), Inches(2.2), col_w - Inches(0.5),
                Inches(2.0), r.get('fond_plus', []),
                size=10, color=TEXT, bullet_color=GREEN,
                line_spacing=1.25, max_lines=MAX_BULLETS + 2)
    # Fond −
    add_text(s, col_x + Inches(0.25), Inches(4.3), col_w - Inches(0.5),
             Inches(0.35),
             '✗ Points faibles', size=12, bold=True, color=RED)
    add_bullets(s, col_x + Inches(0.25), Inches(4.7), col_w - Inches(0.5),
                Inches(2.1), r.get('fond_moins', []),
                size=10, color=TEXT, bullet_color=RED,
                line_spacing=1.25, max_lines=MAX_BULLETS)

    add_footer(s, num, total, f"{r['name']} · Analyse")


def slide_rapport_idees(r, num, total):
    """Slide C : idées éditoriales (titre + pourquoi + comment)."""
    s = prs.slides.add_slide(BLANK)
    sub = f"{r['titre']} · {r['pages']}"
    add_title_bar(s, f"{r['name']} — Idées pour IGENSIA", sub,
                  rank=r['rank'], score=r['score'])

    idees = r.get('idees', [])
    if not idees:
        add_text(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1),
                 'Pas d\'idée éditoriale distinctive retenue : les bonnes pratiques de ce rapport sont déjà couvertes par les autres rapports du benchmark.',
                 size=14, color=TEXT_LIGHT, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
        add_footer(s, num, total, f"{r['name']} · Idées")
        return

    # Layout en colonnes selon nombre d'idées
    n = len(idees)
    if n == 1:
        cols = [(Inches(0.5), Inches(12.3))]
    elif n == 2:
        cols = [(Inches(0.4), Inches(6.3)), (Inches(7.0), Inches(5.95))]
    elif n == 3:
        cols = [(Inches(0.3), Inches(4.3)), (Inches(4.7), Inches(4.3)),
                (Inches(9.1), Inches(4.0))]
    elif n == 4:
        cols = [(Inches(0.3), Inches(3.2)), (Inches(3.6), Inches(3.2)),
                (Inches(6.9), Inches(3.2)), (Inches(10.2), Inches(3.0))]
    else:  # 5+ (OMNES = 6)
        # 2 rangées
        cols_per_row = (n + 1) // 2
        col_w = (Inches(12.7) - Inches(0.15) * (cols_per_row - 1)) / cols_per_row
        cols = []
        for i in range(n):
            row = i // cols_per_row
            col = i % cols_per_row
            x = Inches(0.3) + (col_w + Inches(0.15)) * col
            cols.append((x, col_w, row))

    card_y = Inches(1.25)
    card_h = Inches(5.65)
    # Si > 4 idées, on coupe en 2 rangées
    if n >= 5:
        card_h_each = (card_h - Inches(0.2)) / 2
        for i, c in enumerate(cols):
            x, w, row = c
            y = card_y + (card_h_each + Inches(0.2)) * row
            _render_idee_card(s, x, y, w, card_h_each, idees[i], i + 1,
                              compact=True)
    else:
        for i, (x, w) in enumerate(cols):
            _render_idee_card(s, x, card_y, w, card_h, idees[i], i + 1,
                              compact=(n >= 3))

    add_footer(s, num, total, f"{r['name']} · Idées")


def _render_idee_card(slide, x, y, w, h, idee, num, compact=False):
    """Carte d'une idée : header rose + titre + pourquoi + comment."""
    # Carte fond
    add_rect(slide, x, y, w, h, fill=ROSE_LIGHT, line=ROSE,
             line_w=Emu(9525))
    # Header rose foncé en haut
    add_rect(slide, x, y, w, Inches(0.45), fill=ROSE)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(0.5),
             Inches(0.35),
             f'#{num}', size=12, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1)
    add_text(slide, x + Inches(0.6), y + Inches(0.05), w - Inches(0.75),
             Inches(0.35),
             'IDÉE À REPRENDRE', size=9, bold=True,
             color=WHITE, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1)
    # Titre
    title_y = y + Inches(0.55)
    titre = idee.get('titre', '')
    title_h = Inches(0.85) if not compact else Inches(0.95)
    title_size = 11 if not compact else 10
    add_text(slide, x + Inches(0.2), title_y, w - Inches(0.4), title_h,
             titre, size=title_size, bold=True, color=VIOLET,
             line_spacing=1.2)
    # Pourquoi
    pq_y = title_y + title_h + Inches(0.05)
    add_text(slide, x + Inches(0.2), pq_y, w - Inches(0.4), Inches(0.25),
             'Pourquoi', size=8, bold=True, color=ROSE)
    pq_size = 9 if not compact else 8
    pq_h = (h - (pq_y - y) - Inches(2.0)) if not compact else (h - (pq_y - y) - Inches(1.3))
    add_text(slide, x + Inches(0.2), pq_y + Inches(0.25),
             w - Inches(0.4), pq_h,
             truncate(idee.get('pourquoi', ''), 600 if not compact else 280),
             size=pq_size, color=TEXT, line_spacing=1.25)
    # Comment
    com_y = y + h - (Inches(1.7) if not compact else Inches(1.0))
    add_text(slide, x + Inches(0.2), com_y, w - Inches(0.4), Inches(0.25),
             'Comment', size=8, bold=True, color=ROSE)
    com_size = 9 if not compact else 8
    com_h = (Inches(1.4) if not compact else Inches(0.8))
    add_text(slide, x + Inches(0.2), com_y + Inches(0.25),
             w - Inches(0.4), com_h,
             truncate(idee.get('comment', ''), 500 if not compact else 240),
             size=com_size, color=TEXT, line_spacing=1.25)


# ====================
# CONCLU
# ====================
def slide_conclu(num, total):
    s = prs.slides.add_slide(BLANK)
    add_title_bar(s, 'Top 5 chantiers pour le prochain rapport',
                  'Synthèse des recommandations issues du benchmark')
    actions = [
        ('Tableau d\'indicateurs N / N-1 en fin de rapport',
         'Transforme la photo annuelle en récit de trajectoire. Démarre dès cette année.'),
        ('Connecter chaque action à des cibles ODD numérotées',
         'Au-delà de la mention « ODD 4 », rattacher chaque engagement à 1-2 cibles précises (4.7, 8.5, etc.).'),
        ('Page « Chiffres-clés » en ouverture',
         'Carte d\'identité chiffrée du rapport en une page, 8-10 indicateurs phares.'),
        ('Budget RSE publié en valeur absolue',
         'Le signal de transparence le plus fort du benchmark. Démarre N seul si N-1 non reconstituable.'),
        ('Schéma directeur RSE 5 axes en visuel central',
         'Vue d\'ensemble stratégique avant d\'entrer dans le détail des engagements.'),
    ]
    y = Inches(1.35)
    for i, (title, desc) in enumerate(actions):
        row_h = Inches(1.05)
        # Numéro
        add_rect(s, Inches(0.5), y, Inches(0.7), row_h - Inches(0.1),
                 fill=ROSE)
        add_text(s, Inches(0.5), y, Inches(0.7), row_h - Inches(0.1),
                 f'{i+1}', size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1)
        # Titre + desc
        add_text(s, Inches(1.4), y + Inches(0.05), Inches(11.5),
                 Inches(0.45),
                 title, size=14, bold=True, color=VIOLET,
                 line_spacing=1.15)
        add_text(s, Inches(1.4), y + Inches(0.5), Inches(11.5),
                 Inches(0.45),
                 desc, size=11, color=TEXT, line_spacing=1.25)
        y += row_h
    add_footer(s, num, total, 'Conclusion')


# ====================
# BUILD
# ====================
def main():
    n_per_rapport = 3
    n_intro = 2  # cover + méthodo
    n_igensia = 3  # intro + forces + reco
    n_conclu = 1
    total = n_intro + n_igensia + n_per_rapport * len(RAPPORTS) + n_conclu
    print(f'Total slides : {total}')

    # 1) Cover
    slide_cover()
    # 2) Méthodo
    slide_methodo(2, total)
    # 3) IGENSIA Reference intro
    slide_igensia_intro(3, total)
    # 4) IGENSIA forces
    slide_igensia_forces(4, total)
    # 5) IGENSIA reco
    slide_igensia_reco(5, total)
    # 6) Rapports (tri par rang croissant)
    sorted_r = sorted(RAPPORTS, key=lambda x: x.get('rank', 99))
    cur = 6
    for r in sorted_r:
        print(f"  [{r['rank']:>2}] {r['name']}")
        slide_rapport_visuels(r, cur, total); cur += 1
        slide_rapport_analyse(r, cur, total); cur += 1
        slide_rapport_idees(r, cur, total); cur += 1
    # 7) Conclusion
    slide_conclu(cur, total)

    out = 'Benchmark_RSE_IGENSIA_Rapports_v4.pptx'
    prs.save(out)
    print(f'\n✓ {out} généré ({total} slides)')


if __name__ == '__main__':
    main()
